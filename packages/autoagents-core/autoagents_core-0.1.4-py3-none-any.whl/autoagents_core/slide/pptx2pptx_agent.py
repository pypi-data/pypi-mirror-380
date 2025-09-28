import os
import re
import json
import csv
import tempfile
import requests
import base64
from typing import Optional, Any, Dict, List, Union
from io import BytesIO
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.oxml import parse_xml
from pptx.oxml.ns import qn
from copy import deepcopy


class TextStylePreserver:
    """完整的文本样式保存和恢复工具类"""
    
    @staticmethod
    def capture_complete_style(paragraph):
        """捕获段落的完整样式信息"""
        style_info = {
            'paragraph_level': {},
            'run_level': [],
            'text_frame_level': {}
        }
        
        # 段落级别样式
        style_info['paragraph_level'] = {
            'alignment': paragraph.alignment,
            'level': paragraph.level,
            'space_before': paragraph.space_before,
            'space_after': paragraph.space_after,
            'line_spacing': paragraph.line_spacing,
            'has_bullet': False,
            'bullet_char': None
        }
        
        # 检查是否有项目符号
        try:
            p_element = paragraph._p
            pPr = p_element.pPr
            if pPr is not None:
                # 查找项目符号字符
                buChar_elements = pPr.xpath('.//a:buChar', namespaces={'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'})
                if buChar_elements:
                    style_info['paragraph_level']['has_bullet'] = True
                    char_attr = buChar_elements[0].get('char')
                    if char_attr:
                        style_info['paragraph_level']['bullet_char'] = char_attr
                    else:
                        style_info['paragraph_level']['bullet_char'] = "•"  # 默认符号
                # 检查其他类型的项目符号（如数字编号等）
                elif pPr.xpath('.//a:buFont | .//a:buAutoNum | .//a:buBlip', namespaces={'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}):
                    style_info['paragraph_level']['has_bullet'] = True
                    style_info['paragraph_level']['bullet_char'] = "•"  # 默认符号
        except Exception:
            pass
        
        # 获取段落的XML以保存更多属性
        try:
            # 保存段落属性的XML片段
            style_info['paragraph_level']['xml_props'] = p_element.xml if hasattr(p_element, 'xml') else None
        except:
            pass
        
        # Run级别样式（每个run）
        for run in paragraph.runs:
            # 处理可能为None的属性
            font_bold = run.font.bold
            if font_bold is None:
                font_bold = False  # 默认不bold
            
            font_italic = run.font.italic  
            if font_italic is None:
                font_italic = False  # 默认不italic
                
            font_underline = run.font.underline
            if font_underline is None:
                font_underline = False  # 默认无下划线
            
            run_style = {
                'text': run.text,
                'font_name': run.font.name,
                'font_size': run.font.size,
                'font_bold': font_bold,
                'font_italic': font_italic,
                'font_underline': font_underline,
                'font_color_rgb': None,
                'font_color_theme': None,
                'hyperlink': None,
                'font_element_xml': None  # 保存字体元素的XML
            }
            
            # 字体颜色处理
            try:
                if run.font.color.rgb:
                    run_style['font_color_rgb'] = run.font.color.rgb
                elif run.font.color.theme_color:
                    run_style['font_color_theme'] = run.font.color.theme_color
            except AttributeError:
                pass
            
            # 尝试保存字体的原始XML以确保完整性
            try:
                if hasattr(run, '_r'):
                    run_style['font_element_xml'] = run._r.xml
            except:
                pass
            
            # 超链接处理
            try:
                if hasattr(run, '_r') and run._r.get('hlinkClick'):
                    run_style['hyperlink'] = run._r.get('hlinkClick')
            except:
                pass
                
            style_info['run_level'].append(run_style)
        
        return style_info
    
    @staticmethod
    def capture_text_frame_style(text_frame):
        """捕获文本框级别的样式"""
        return {
            'vertical_anchor': text_frame.vertical_anchor,
            'margin_left': text_frame.margin_left,
            'margin_right': text_frame.margin_right,
            'margin_top': text_frame.margin_top,
            'margin_bottom': text_frame.margin_bottom,
            'word_wrap': text_frame.word_wrap,
            'auto_size': text_frame.auto_size
        }
    
    @staticmethod
    def apply_style_to_new_text(paragraph, style_info, new_text):
        """将保存的样式应用到新文本上，完全保留格式"""
        
        # 应用段落级别样式
        para_style = style_info['paragraph_level']
        if para_style.get('alignment') is not None:
            paragraph.alignment = para_style['alignment']
        if para_style.get('level') is not None:
            paragraph.level = para_style['level']
        if para_style.get('space_before') is not None:
            paragraph.space_before = para_style['space_before']
        if para_style.get('space_after') is not None:
            paragraph.space_after = para_style['space_after']
        if para_style.get('line_spacing') is not None:
            paragraph.line_spacing = para_style['line_spacing']
        
        # 恢复项目符号
        if para_style.get('has_bullet', False):
            bullet_char = para_style.get('bullet_char', "•")
            enable_bullet(paragraph, bullet_char)
        
        # 保留第一个run的样式信息，然后替换文本
        if paragraph.runs and style_info['run_level']:
            original_run_style = style_info['run_level'][0]
            
            # 清除所有现有runs（除了第一个）
            while len(paragraph.runs) > 1:
                run = paragraph.runs[-1]
                run._r.getparent().remove(run._r)
            
            # 获取第一个run并设置新文本
            if paragraph.runs:
                run = paragraph.runs[0]
                run.text = new_text
                
                # 强化字体样式应用 - 确保字体名称被正确设置
                try:
                    if original_run_style.get('font_name'):
                        original_font_name = original_run_style['font_name']
                        run.font.name = original_font_name
                        
                        # 验证字体是否设置成功，如果不匹配则尝试XML级别设置
                        if run.font.name != original_font_name:
                            try:
                                # 直接操作底层XML元素
                                if hasattr(run, '_r'):
                                    r_element = run._r
                                    # 获取或创建rPr元素
                                    rPr = r_element.get_or_add_rPr()
                                    
                                    # 移除现有的latin字体设置（如果有）
                                    for latin in rPr.xpath('.//a:latin', namespaces={'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}):
                                        rPr.remove(latin)
                                    
                                    # 添加新的latin字体设置
                                    from pptx.oxml import parse_xml
                                    latin_font = parse_xml(f'<a:latin xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" typeface="{original_font_name}"/>')
                                    rPr.insert(0, latin_font)
                            except Exception:
                                pass  # 字体设置失败时静默处理
                except Exception:
                    pass  # 字体名称设置异常时静默处理
                
                try:
                    if original_run_style.get('font_size'):
                        run.font.size = original_run_style['font_size']
                except Exception as e:
                    print(f"⚠️ 字体大小设置异常: {e}")
                
                # 明确设置布尔属性
                try:
                    run.font.bold = original_run_style.get('font_bold', False)
                except:
                    pass
                try:
                    run.font.italic = original_run_style.get('font_italic', False)
                except:
                    pass
                try:
                    run.font.underline = original_run_style.get('font_underline', False)
                except:
                    pass
                
                # 应用字体颜色
                try:
                    if original_run_style.get('font_color_rgb'):
                        run.font.color.rgb = original_run_style['font_color_rgb']
                    elif original_run_style.get('font_color_theme'):
                        run.font.color.theme_color = original_run_style['font_color_theme']
                except Exception as e:
                    print(f"⚠️ 字体颜色设置异常: {e}")
            else:
                # 如果没有现有runs，创建一个新的run并应用样式
                paragraph.text = new_text
                if paragraph.runs:
                    run = paragraph.runs[0]
                    # 强化字体样式应用（备用情况）
                    try:
                        if original_run_style.get('font_name'):
                            original_font_name = original_run_style['font_name']
                            run.font.name = original_font_name
                            # 如果字体名称不匹配，尝试XML级别设置
                            if run.font.name != original_font_name:
                                try:
                                    # 直接操作底层XML元素
                                    if hasattr(run, '_r'):
                                        r_element = run._r
                                        # 获取或创建rPr元素
                                        rPr = r_element.get_or_add_rPr()
                                        
                                        # 移除现有的latin字体设置（如果有）
                                        for latin in rPr.xpath('.//a:latin', namespaces={'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}):
                                            rPr.remove(latin)
                                        
                                        # 添加新的latin字体设置
                                        from pptx.oxml import parse_xml
                                        latin_font = parse_xml(f'<a:latin xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" typeface="{original_font_name}"/>')
                                        rPr.insert(0, latin_font)
                                except Exception:
                                    pass  # 字体设置失败时静默处理
                    except Exception:
                        pass  # 字体名称设置异常时静默处理
                    
                    try:
                        if original_run_style.get('font_size'):
                            run.font.size = original_run_style['font_size']
                    except:
                        pass
                    try:
                        run.font.bold = original_run_style.get('font_bold', False)
                    except:
                        pass
                    try:
                        run.font.italic = original_run_style.get('font_italic', False)
                    except:
                        pass
                    try:
                        run.font.underline = original_run_style.get('font_underline', False)
                    except:
                        pass
                    try:
                        if original_run_style.get('font_color_rgb'):
                            run.font.color.rgb = original_run_style['font_color_rgb']
                        elif original_run_style.get('font_color_theme'):
                            run.font.color.theme_color = original_run_style['font_color_theme']
                    except:
                        pass
        else:
            # 如果没有样式信息，使用简单的文本替换
            paragraph.text = new_text
    
    @staticmethod
    def apply_text_frame_style(text_frame, style_info):
        """应用文本框级别样式"""
        if style_info.get('vertical_anchor') is not None:
            text_frame.vertical_anchor = style_info['vertical_anchor']
        if style_info.get('margin_left') is not None:
            text_frame.margin_left = style_info['margin_left']
        if style_info.get('margin_right') is not None:
            text_frame.margin_right = style_info['margin_right']
        if style_info.get('margin_top') is not None:
            text_frame.margin_top = style_info['margin_top']
        if style_info.get('margin_bottom') is not None:
            text_frame.margin_bottom = style_info['margin_bottom']
        if style_info.get('word_wrap') is not None:
            text_frame.word_wrap = style_info['word_wrap']
        if style_info.get('auto_size') is not None:
            text_frame.auto_size = style_info['auto_size']


def replace_text_preserve_format(text_frame, new_text):
    """替换文本并完全保留格式的核心函数"""
    if not text_frame.paragraphs:
        return
    
    # 捕获文本框级别样式
    tf_style = TextStylePreserver.capture_text_frame_style(text_frame)
    
    # 检查新文本是否包含多行
    new_lines = new_text.split('\n')
    original_paragraphs = list(text_frame.paragraphs)
    
    # 捕获所有现有段落的样式
    para_styles = []
    for para in original_paragraphs:
        para_styles.append(TextStylePreserver.capture_complete_style(para))
    
    # 如果新文本只有一行，使用原来的简单逻辑
    if len(new_lines) == 1:
        # 捕获第一个段落的样式作为模板
        first_para = text_frame.paragraphs[0]
        para_style = para_styles[0] if para_styles else TextStylePreserver.capture_complete_style(first_para)
        
        # 删除所有现有段落（除了第一个）
        while len(text_frame.paragraphs) > 1:
            p = text_frame.paragraphs[-1]
            text_frame._element.remove(p._p)
        
        # 在第一个段落应用新文本和样式
        TextStylePreserver.apply_style_to_new_text(first_para, para_style, new_text)
    else:
        # 多行文本，需要保留每行的段落格式
        # 首先处理第一行
        if para_styles:
            first_para = text_frame.paragraphs[0]
            TextStylePreserver.apply_style_to_new_text(first_para, para_styles[0], new_lines[0])
        
        # 处理剩余行
        for i, line in enumerate(new_lines[1:], 1):
            # 决定使用哪个样式模板
            if i < len(para_styles):
                # 使用对应的原始段落样式
                style_template = para_styles[i]
            elif para_styles:
                # 使用最后一个段落的样式作为模板
                style_template = para_styles[-1]
            else:
                # 使用第一个段落的样式作为模板
                style_template = para_styles[0] if para_styles else TextStylePreserver.capture_complete_style(text_frame.paragraphs[0])
            
            # 添加新段落或使用现有段落
            if i < len(text_frame.paragraphs):
                # 使用现有段落
                para = text_frame.paragraphs[i]
            else:
                # 添加新段落
                para = text_frame.add_paragraph()
            
            # 应用样式和文本
            TextStylePreserver.apply_style_to_new_text(para, style_template, line)
        
        # 删除多余的段落（如果新文本行数少于原段落数）
        while len(text_frame.paragraphs) > len(new_lines):
            p = text_frame.paragraphs[-1]
            text_frame._element.remove(p._p)
    
    # 恢复文本框级别样式
    TextStylePreserver.apply_text_frame_style(text_frame, tf_style)


def process_list_preserve_format(text_frame, list_data):
    """处理列表数据并完全保留格式"""
    if not text_frame.paragraphs or not list_data:
        return
        
    # 捕获文本框级别样式
    tf_style = TextStylePreserver.capture_text_frame_style(text_frame)
    
    # 捕获第一个段落的样式作为模板
    first_para = text_frame.paragraphs[0]
    para_style = TextStylePreserver.capture_complete_style(first_para)
    
    # 删除所有现有段落（除了第一个）
    while len(text_frame.paragraphs) > 1:
        p = text_frame.paragraphs[-1]
        text_frame._element.remove(p._p)
    
    # 处理第一个列表项
    if list_data:
        first_item = str(list_data[0])
        TextStylePreserver.apply_style_to_new_text(first_para, para_style, first_item)
        enable_bullet(first_para)
        
        # 处理剩余的列表项
        for item in list_data[1:]:
            new_para = text_frame.add_paragraph()
            TextStylePreserver.apply_style_to_new_text(new_para, para_style, str(item))
            enable_bullet(new_para)
    
    # 恢复文本框级别样式
    TextStylePreserver.apply_text_frame_style(text_frame, tf_style)


def parse_markdown_text_preserve_format(text_frame, markdown_text):
    """
    解析Markdown文本并完全保留原有格式
    支持：
    - * bullet points
    - **粗体**
    - *斜体*
    - `代码`
    - # 标题
    """
    if not text_frame.paragraphs:
        return
        
    # 捕获文本框级别样式
    tf_style = TextStylePreserver.capture_text_frame_style(text_frame)
    
    # 捕获第一个段落的样式作为模板
    first_para = text_frame.paragraphs[0]
    para_style = TextStylePreserver.capture_complete_style(first_para)
    
    # 删除所有现有段落（除了第一个）
    while len(text_frame.paragraphs) > 1:
        p = text_frame.paragraphs[-1]
        text_frame._element.remove(p._p)
    
    lines = markdown_text.split('\n')
    first_line_processed = False
    
    for line in lines:
        line = line.strip()
        if not line:
            continue
            
        # 选择要处理的段落
        if not first_line_processed:
            p = first_para
            first_line_processed = True
        else:
            p = text_frame.add_paragraph()
        
        # 处理标题 (# ## ###)
        if line.startswith('#'):
            level = 0
            while level < len(line) and line[level] == '#':
                level += 1
            title_text = line[level:].strip()
            
            # 应用原有样式，但文本是标题
            TextStylePreserver.apply_style_to_new_text(p, para_style, title_text)
            # 可以考虑适当增大字体大小表示标题，但保留其他格式
            if p.runs and para_style['run_level']:
                run = p.runs[0]
                original_size = para_style['run_level'][0].get('font_size')
                if original_size:
                    # 根据标题级别适当增加字体大小
                    size_increase = max(0, (4 - level) * 2)
                    try:
                        run.font.size = Pt(original_size.pt + size_increase)
                    except:
                        pass
                run.font.bold = True
            continue
        
        # 处理bullet points
        if line.startswith('* ') or line.startswith('- '):
            bullet_text = line[2:].strip()
            
            # 检查是否包含内联格式
            if any(marker in bullet_text for marker in ['**', '*', '`']):
                apply_inline_formatting_preserve_format(p, bullet_text, para_style)
            else:
                TextStylePreserver.apply_style_to_new_text(p, para_style, bullet_text)
            
            enable_bullet(p)
            continue
        
        # 处理普通文本（可能包含内联格式）
        if any(marker in line for marker in ['**', '*', '`']):
            apply_inline_formatting_preserve_format(p, line, para_style)
        else:
            TextStylePreserver.apply_style_to_new_text(p, para_style, line)
    
    # 如果没有处理任何行，至少清空第一个段落
    if not first_line_processed:
        TextStylePreserver.apply_style_to_new_text(first_para, para_style, "")
    
    # 恢复文本框级别样式
    TextStylePreserver.apply_text_frame_style(text_frame, tf_style)


def apply_inline_formatting_preserve_format(paragraph, text, base_style):
    """
    应用内联格式：粗体、斜体、代码，并保留基础样式
    
    对于复杂的内联格式，暂时先设置普通文本，保留原有样式
    TODO: 在未来版本中可以改进为完全支持内联格式的样式保留
    """
    # 暂时简化处理，直接应用文本并保留样式
    TextStylePreserver.apply_style_to_new_text(paragraph, base_style, text)


def get_jwt_token_api(
    personal_auth_key: str,
    personal_auth_secret: str,
    base_url: str = "https://uat.agentspro.cn",
) -> str:
    """
    获取 autoagents_core AI 平台的 JWT 认证令牌，用户级认证，用于后续的 API 调用认证。
    JWT token 具有时效性，30天过期后需要重新获取。
    
    Args:
        agent_id (str): Agent 的唯一标识符，用于调用Agent对话
            - 获取方式：Agent详情页 - 分享 - API
            
        personal_auth_key (str): 认证密钥
            - 获取方式：右上角 - 个人密钥
            
        personal_auth_secret (str): 认证密钥
            - 获取方式：右上角 - 个人密钥

        base_url (str, optional): API 服务基础地址
            - 默认值: "https://uat.agentspro.cn"
            - 测试环境: "https://uat.agentspro.cn"  
            - 生产环境: "https://agentspro.cn"
            - 私有部署时可指定自定义地址
            
    Returns:
        str: JWT 认证令牌            
    """
    
    headers = {
        "Authorization": f"Bearer {personal_auth_key}.{personal_auth_secret}",
        "Content-Type": "application/json"
    }

    url = f"{base_url}/openapi/user/auth"
    response = requests.get(url, headers=headers)
    return response.json()["data"]["token"]

def extract_json(text: str | None = None):
    """从AI响应中提取JSON内容，处理各种格式情况"""
    if not text:
        return None

    # 1. 先尝试提取```json```代码块中的JSON
    json_pattern = r'```json\s*(.*?)\s*```'
    match = re.search(json_pattern, text, re.DOTALL | re.IGNORECASE)
    if match:
        json_str = match.group(1).strip()
        try:
            return json.loads(json_str)
        except json.JSONDecodeError:
            pass

    # 2. 尝试提取纯JSON代码块（无语言标识）
    code_block_pattern = r'```\s*([\s\S]*?)\s*```'
    match = re.search(code_block_pattern, text, re.DOTALL)
    if match:
        potential_json = match.group(1).strip()
        try:
            return json.loads(potential_json)
        except json.JSONDecodeError:
            pass
    
    # 3. 尝试直接查找JSON对象（以{开头，以}结尾）
    json_start = text.find('{')
    json_end = text.rfind('}') + 1
    if json_start != -1 and json_end > json_start:
        potential_json = text[json_start:json_end]
        try:
            return json.loads(potential_json)
        except json.JSONDecodeError:
            pass

    return None


def convert_csv_to_json_list(csv_file_path: str):
    """
    读取CSV文件并转换为json列表
    """
    try:
        if not os.path.exists(csv_file_path):
            print(f"CSV文件不存在: {csv_file_path}")
            return []
        
        data = []
        with open(csv_file_path, 'r', encoding='utf-8') as file:
            reader = csv.DictReader(file)
            for row in reader:
                # 转换数值类型
                converted_row = {}
                for key, value in row.items():
                    # 尝试转换为数字
                    try:
                        if '.' in value:
                            converted_row[key] = float(value)
                        else:
                            converted_row[key] = int(value)
                    except (ValueError, TypeError):
                        # 如果转换失败，保持原始字符串
                        converted_row[key] = value
                
                data.append(converted_row)
        
        print(f"成功读取CSV: {csv_file_path}, {len(data)} 条记录")
        return data
    except Exception as e:
        print(f"读取CSV文件失败: {csv_file_path}, 错误: {e}")
        return []


def parse_markdown_text(text_frame, markdown_text, font_size=14, preserve_style=True):
    """
    解析Markdown文本并应用到PowerPoint文本框
    
    DEPRECATED: 这个函数为了兼容性保留，建议使用 parse_markdown_text_preserve_format
    现在内部使用新的格式保留逻辑，完全保留原有格式。
    
    Args:
        preserve_style: 是否保留原有样式，默认True保留模板样式（总是True）
        font_size: 废弃参数，不再使用
    """
    # 现在总是使用新的格式保留函数
    parse_markdown_text_preserve_format(text_frame, markdown_text)


def apply_inline_formatting(paragraph, text, preserve_style=True):
    """
    应用内联格式：粗体、斜体、代码
    
    DEPRECATED: 这个函数为了兼容性保留，建议使用 apply_inline_formatting_preserve_format
    现在内部使用新的格式保留逻辑。
    
    Args:
        preserve_style: 是否保留原有样式，默认True（总是True）
    """
    # 首先捕获当前段落的样式
    base_style = TextStylePreserver.capture_complete_style(paragraph)
    
    # 使用新的格式保留函数
    apply_inline_formatting_preserve_format(paragraph, text, base_style)


def enable_bullet(paragraph, bullet_char="•"):
    """为段落启用项目符号"""
    p = paragraph._p
    pPr = p.get_or_add_pPr()
    buChar = parse_xml(f'<a:buChar xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" char="{bullet_char}"/>')
    pPr.insert(0, buChar)


def fill_existing_table(table, data, font_size=12):
    """
    将 data 填充到 pptx 表格中，第二行作为模板行。
    支持字段替换：[name], [count], [=count*price], [@picture]
    """
    from pptx.util import Pt
    import re

    def eval_formula(expr, context):
        try:
            return str(eval(expr, {}, context))
        except:
            return expr

    # 数据类型检查和转换
    if not isinstance(data, list):
        print(f"错误: 表格数据必须是列表格式，当前类型: {type(data)}")
        return
    
    if not data:
        print("警告: 表格数据为空")
        return
    
    # 确保列表中的每个元素都是字典
    for i, item in enumerate(data):
        if not isinstance(item, dict):
            print(f"错误: 表格数据第{i}项不是字典格式，当前类型: {type(item)}")
            return

    # 查找包含占位符的行作为模板行
    template_row_idx = None
    for i, row in enumerate(table.rows):
        for cell in row.cells:
            if '[' in cell.text and ']' in cell.text:
                template_row_idx = i
                break
        if template_row_idx is not None:
            break
    
    if template_row_idx is None:
        return
        
    template_row = table.rows[template_row_idx]
    col_templates = [cell.text for cell in template_row.cells]
    
    # 获取表格的底层XML对象（在循环外定义）
    tbl = table._tbl

    for row_data in data:
        # 使用底层XML操作添加新行
        new_tr = parse_xml(template_row._tr.xml)
        tbl.append(new_tr)
        
        # 获取新添加的行
        new_row_idx = len(table.rows) - 1
        row = table.rows[new_row_idx]
        
        for j, tmpl in enumerate(col_templates):
            text = tmpl
            
            # 处理图片占位符 [@picture]
            if re.search(r'\[@\w+\]', text):
                # 找到图片占位符
                img_match = re.search(r'\[@(\w+)\]', text)
                if img_match:
                    img_key = img_match.group(1)
                    if img_key in row_data:
                        # 清空单元格文本
                        row.cells[j].text = ""
                        # 这里可以添加图片到单元格的逻辑
                        # 目前先显示图片文件名
                        row.cells[j].text = f"图片: {row_data[img_key]}"
                    continue
            
            # 字段替换
            for key, val in row_data.items():
                text = text.replace(f"[{key}]", str(val))
            
            # 表达式处理（如 [=count*price]）
            match = re.findall(r"\[=([^\]]+)\]", text)
            for m in match:
                result = eval_formula(m, row_data)
                text = text.replace(f"[={m}]", result)
            
            # 检查是否包含Markdown格式
            if any(marker in text for marker in ['*', '#', '`']):
                # 使用Markdown解析，完全保留格式
                parse_markdown_text_preserve_format(row.cells[j].text_frame, text)
            else:
                # 普通文本 - 使用新的格式保留函数
                replace_text_preserve_format(row.cells[j].text_frame, text)

    # 删除模板行
    tbl.remove(template_row._tr)


def find_nearest_table(placeholder_shape, all_tables):
    """
    根据最近距离原则找到对应的表格
    """
    if not all_tables:
        return None
    
    placeholder_pos = (placeholder_shape.left, placeholder_shape.top)
    
    def calculate_distance(table_shape):
        table_pos = (table_shape.left, table_shape.top)
        return ((placeholder_pos[0] - table_pos[0]) ** 2 + 
                (placeholder_pos[1] - table_pos[1]) ** 2) ** 0.5
    
    return min(all_tables, key=calculate_distance)


def create_temp_file(file_content, file_extension='.pptx'):
    """
    创建临时文件并写入内容
    
    Args:
        file_content: 文件内容（字节或字符串）
        file_extension: 文件扩展名，默认为.pptx
    
    Returns:
        str: 临时文件路径
    """
    temp_fd, temp_path = tempfile.mkstemp(suffix=file_extension)
    try:
        with os.fdopen(temp_fd, 'wb' if isinstance(file_content, bytes) else 'w') as tmp_file:
            tmp_file.write(file_content)
        return temp_path
    except Exception as e:
        os.unlink(temp_path)
        raise e


def cleanup_temp_file(temp_path):
    """
    清理临时文件
    
    Args:
        temp_path: 临时文件路径
    """
    try:
        if os.path.exists(temp_path):
            os.unlink(temp_path)
    except Exception:
        pass  # 忽略清理错误


def download_image(url: str) -> Optional[str]:
    """
    下载远程图片到临时文件
    
    Args:
        url: 图片URL
    
    Returns:
        str: 临时文件路径，如果下载失败返回None
    """
    
    try:
        print(f"开始下载图片: {url}")
        
        # 添加User-Agent头，避免被某些网站拒绝
        headers = {
            'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36'
        }
        
        response = requests.get(url, stream=True, headers=headers, timeout=30)
        response.raise_for_status()
        
        # 验证content-type是否为图片
        content_type = response.headers.get('content-type', '').lower()
        print(f"图片URL Content-Type: {content_type}")
        
        valid_image_types = [
            'image/jpeg', 'image/jpg', 'image/png', 'image/gif', 
            'image/bmp', 'image/webp', 'image/tiff', 'image/svg+xml'
        ]
        
        if not any(img_type in content_type for img_type in valid_image_types):
            print(f"❌ 跳过非图片内容: {url}, content-type: {content_type}")
            return None
        
        print(f"✅ 确认为图片内容: {content_type}")
        
        # 创建临时文件
        suffix = '.jpg'  # 默认后缀
        if 'image/png' in content_type:
            suffix = '.png'
        elif 'image/gif' in content_type:
            suffix = '.gif'
        elif 'image/webp' in content_type:
            suffix = '.webp'
        elif 'image/bmp' in content_type:
            suffix = '.bmp'
        elif 'image/tiff' in content_type:
            suffix = '.tiff'
        elif 'image/svg' in content_type:
            suffix = '.svg'
        
        temp_file = tempfile.NamedTemporaryFile(delete=False, suffix=suffix)
        
        # 下载图片内容
        total_size = 0
        for chunk in response.iter_content(chunk_size=8192):
            temp_file.write(chunk)
            total_size += len(chunk)
        
        temp_file.close()
        print(f"✅ 图片下载完成: {temp_file.name} (大小: {total_size} 字节)")
        return temp_file.name
        
    except Exception as e:
        print(f"❌ 下载图片失败: {url}")
        print(f"   错误详情: {type(e).__name__}: {e}")
        return None


def download_template(url: str) -> Optional[str]:
    """
    下载远程模板文件到临时文件
    
    Args:
        url: 模板文件URL
    
    Returns:
        str: 临时文件路径，如果下载失败返回None
    """
    
    try:
        response = requests.get(url, stream=True)
        response.raise_for_status()
        
        # 验证content-type是否为PowerPoint文件
        content_type = response.headers.get('content-type', '').lower()
        valid_ppt_types = [
            'application/vnd.openxmlformats-officedocument.presentationml.presentation',
            'application/vnd.ms-powerpoint',
            'application/octet-stream'  # 有些服务器可能返回这个通用类型
        ]
        
        # 从URL路径获取文件扩展名
        url_path = url.split('?')[0].lower()  # 去掉查询参数
        file_extension = None
        if url_path.endswith(('.pptx', '.ppt')):
            file_extension = '.pptx' if url_path.endswith('.pptx') else '.ppt'
        
        # 如果content-type不匹配但URL有正确的扩展名，仍然尝试下载
        if not any(ppt_type in content_type for ppt_type in valid_ppt_types) and not file_extension:
            print(f"跳过非PowerPoint文件: {url}, content-type: {content_type}")
            return None
        
        # 创建临时文件，优先使用URL中的扩展名
        suffix = file_extension or '.pptx'  # 默认使用.pptx
        temp_file = tempfile.NamedTemporaryFile(delete=False, suffix=suffix)
        
        # 下载文件内容
        for chunk in response.iter_content(chunk_size=8192):
            temp_file.write(chunk)
        
        temp_file.close()
        print(f"成功下载模板文件: {url} -> {temp_file.name}")
        return temp_file.name
        
    except Exception as e:
        print(f"下载模板文件失败: {url}, 错误: {e}")
        return None


def normalize_data_format(data: dict) -> dict:
    """
    标准化数据格式，自动处理带外层包装的数据
    
    支持的输入格式：
    1. 直接数据格式: {"user": {"name": "frank"}, "score": 95}
    2. 带result包装: {"result": {"user": {"name": "frank"}, "score": 95}}
    
    Args:
        data: 输入的数据字典
    
    Returns:
        dict: 标准化后的数据字典
    """
    if not isinstance(data, dict):
        return data
    
    # 检查是否有外层包装
    wrapper_keys = ['result']
    
    for wrapper_key in wrapper_keys:
        if wrapper_key in data and len(data) == 1:
            # 如果只有一个key且是包装key，提取内部数据
            inner_data = data[wrapper_key]
            if isinstance(inner_data, dict):
                print(f"🔍 检测到外层包装 '{wrapper_key}'，自动提取内部数据")
                return inner_data
    
    # 如果没有外层包装，直接返回原数据
    return data


def get_value_by_key(data: dict, key: str) -> Any:
    """
    根据key从数据中获取值，支持点号分隔的嵌套路径和数组索引
    支持格式：
    - 简单key: title, content, table_data
    - 嵌套路径: user.nickname, user.dad.nickname, user.hobbies
    - 数组索引: page[0], items[1]
    - 复合路径: page[0].title, user.hobbies[2], page[1].sections[0].content
    """
    import re
    
    try:
        # 如果key不包含点号和方括号，直接返回
        if '.' not in key and '[' not in key:
            return data.get(key, None)
        
        # 处理复合路径，支持数组索引
        current = data
        
        # 分割路径，同时处理数组索引
        # 例如: "page[0].title" -> ["page[0]", "title"]
        # 例如: "page[1].sections[0].content" -> ["page[1]", "sections[0]", "content"]
        path_parts = []
        current_part = ""
        bracket_depth = 0
        
        for char in key:
            if char == '[':
                bracket_depth += 1
                current_part += char
            elif char == ']':
                bracket_depth -= 1
                current_part += char
            elif char == '.' and bracket_depth == 0:
                if current_part:
                    path_parts.append(current_part)
                    current_part = ""
            else:
                current_part += char
        
        if current_part:
            path_parts.append(current_part)
        
        # 逐级访问每个路径部分
        for part in path_parts:
            # 检查是否包含数组索引
            if '[' in part and ']' in part:
                # 解析数组索引: "page[0]" -> ("page", 0)
                match = re.match(r'^([^[]+)\[(\d+)\]$', part)
                if match:
                    array_name = match.group(1)
                    array_index = int(match.group(2))
                    
                    # 先访问数组
                    if isinstance(current, dict):
                        current = current.get(array_name, None)
                        if current is None:
                            return None
                    else:
                        return None
                    
                    # 再访问索引
                    if isinstance(current, list) and 0 <= array_index < len(current):
                        current = current[array_index]
                    else:
                        return None
                else:
                    # 格式不正确的数组索引
                    return None
            else:
                # 普通key访问
                if isinstance(current, dict):
                    current = current.get(part, None)
                    if current is None:
                        return None
                else:
                    return None
        
        return current
    except Exception as e:
        print(f"获取数据错误: {key}, 错误: {e}")
        return None


def replace_mixed_placeholders(text: str, data: dict) -> str:
    """
    替换文本中的混合占位符
    支持格式: "欢迎 {{name}}，今天是 {{date}}"
    
    Args:
        text: 包含占位符的文本
        data: 数据字典
    
    Returns:
        str: 替换后的文本
    """
    # 使用正则表达式找到所有占位符
    placeholder_pattern = r'\{\{([^}]+)\}\}'
    
    def replace_placeholder(match):
        placeholder_content = match.group(1).strip()
        
        # 判断类型前缀
        if placeholder_content.startswith("@"):
            # 图片占位符在混合文本中不支持，返回原文
            return match.group(0)
        elif placeholder_content.startswith("#"):
            # 表格占位符在混合文本中不支持，返回原文
            return match.group(0)
        else:
            # 普通文本占位符
            value = get_value_by_key(data, placeholder_content)
            if value is not None:
                return str(value)
            else:
                # 如果找不到值，保留原占位符
                return match.group(0)
    
    # 替换所有占位符
    result = re.sub(placeholder_pattern, replace_placeholder, text)
    return result


def is_pure_placeholder(text: str) -> Optional[str]:
    """
    检查文本是否为纯占位符（只包含一个占位符且无其他文字）
    
    Args:
        text: 要检查的文本
    
    Returns:
        str: 如果是纯占位符，返回占位符内容；否则返回None
    """
    text = text.strip()
    if text.startswith("{{") and text.endswith("}}") and text.count("{{") == 1:
        return text[2:-2].strip()
    return None


# Uploader related functions (extracted from utils/uploader.py)
def create_file_like(file_input, filename: Optional[str] = None):
    """创建类文件对象"""
    # 处理不同类型的输入
    if isinstance(file_input, str):
        # 文件路径
        with open(file_input, "rb") as f:
            file_content = f.read()
        
        file_like = BytesIO(file_content)
        file_like.name = file_input.split("/")[-1]
        return file_like
        
    elif isinstance(file_input, bytes):
        # 原始字节数据
        file_like = BytesIO(file_input)
        file_like.name = filename or "uploaded_file"
        return file_like
        
    elif isinstance(file_input, BytesIO):
        # 已经是 BytesIO 对象
        if not hasattr(file_input, 'name') or not file_input.name:
            file_input.name = filename or "uploaded_file"
        return file_input
        
    elif hasattr(file_input, 'read'):
        # 文件对象或类文件对象
        try:
            # 尝试读取内容
            content = file_input.read()
            if isinstance(content, str):
                content = content.encode('utf-8')
            
            file_like = BytesIO(content)
            
            # 确定文件名的优先级
            if filename:
                file_like.name = filename
            elif hasattr(file_input, 'filename') and file_input.filename:
                file_like.name = file_input.filename
            elif hasattr(file_input, 'name') and file_input.name:
                file_like.name = os.path.basename(file_input.name)
            else:
                file_like.name = "uploaded_file"
                
            return file_like
            
        except Exception as e:
            raise ValueError(f"无法读取文件对象: {str(e)}")
            
    else:
        raise TypeError(f"不支持的文件输入类型: {type(file_input)}")


class SimpleFileUploader:
    """简化的文件上传器"""
    
    def __init__(self, personal_auth_key: str, personal_auth_secret: str, base_url: str = "https://uat.agentspro.cn"):
        self.jwt_token = get_jwt_token_api(personal_auth_key, personal_auth_secret, base_url)
        self.base_url = base_url
        self.headers = {
            "Authorization": f"Bearer {self.jwt_token}"
        }

    def upload(self, file, filename: str = "uploaded_file") -> Dict:
        """上传文件"""
        import mimetypes
        
        url = f"{self.base_url}/api/fs/upload"
        
        # 根据文件扩展名自动检测MIME类型
        mime_type, _ = mimetypes.guess_type(filename)
        if mime_type is None:
            mime_type = 'application/octet-stream'  # 默认类型
        
        print(f"Debug: 上传文件 {filename}, 检测到MIME类型: {mime_type}")
        
        files = [
            ('file', (filename, file, mime_type))
        ]
        
        payload = {}
        
        try:
            response = requests.post(url, headers=self.headers, data=payload, files=files, timeout=30)
            
            if response.status_code == 200:
                try:
                    result = response.json()
                    if result.get('code') == 1:  # 成功
                        file_id = result["data"]
                        return {
                            "fileId": file_id,
                            "fileName": filename,
                            "fileType": mime_type,
                            "fileUrl": "",  # 当前API不返回URL
                            "success": True
                        }
                    else:  # 失败
                        error_msg = result.get('msg', '未知错误')
                        raise Exception(f"API返回错误: {error_msg}")
                        
                except Exception as e:
                    # 如果不是JSON响应，返回错误信息字典
                    print(f"Debug: 非JSON响应，返回原始文本: {response.text}")
                    return {
                        "fileId": "",
                        "fileName": filename,
                        "fileType": mime_type,
                        "fileUrl": "",
                        "success": False,
                        "error": response.text.strip()
                    }
            else:
                raise Exception(f"Upload failed: {response.status_code} - {response.text}")
        except Exception as e:
            raise Exception(f"File upload error: {str(e)}")


class PPTX2PPTXAgent:
    """
    独立的PPT填充代理，支持嵌套JSON数据结构填充PowerPoint模板
    
    支持的占位符格式：
    - {{key}} : 文本占位符，对应data['key']
    - {{@key}}: 图片占位符，对应data['key']（图片路径或URL）
    - {{#key}}: 表格占位符，对应data['key']（列表数据）
    
    支持的数据访问路径：
    - 简单路径: {{title}} → data['title']
    - 嵌套路径: {{user.nickname}} → data['user']['nickname'] 
    - 深度嵌套: {{user.dad.nickname}} → data['user']['dad']['nickname']
    - 数组访问: {{user.hobbies}} → data['user']['hobbies'] (整个数组)
    - 数组索引: {{page[0]}} → data['page'][0]
    - 复合索引: {{page[0].title}} → data['page'][0]['title']
    - 深度索引: {{page[1].sections[0].content}} → data['page'][1]['sections'][0]['content']
    
    支持的数据格式：
    1. 直接数据格式: {"user": {"name": "frank"}, "score": 95}
    2. 带外层包装格式: 
       - {"result": {"user": {"name": "frank"}, "score": 95}}
       - {"data": {"user": {"name": "frank"}, "score": 95}}
       - {"payload": {"user": {"name": "frank"}, "score": 95}}
       - {"content": {"user": {"name": "frank"}, "score": 95}}
    
    其他特性：
    - 自动检测并提取外层包装数据
    - 完全保留文本格式（颜色、字体、对齐方式等）
    - 支持Markdown格式解析
    - 支持嵌套字典结构和点号分隔的路径访问
    - 支持数组索引访问 (如 page[0].title, user.hobbies[1])
    - 支持动态重排序：根据数据中的order信息重新排列幻灯片
    """
    
    def __init__(self):
        pass
    
    def _is_safe_shape(self, shape, verbose: bool = False):
        """
        安全检查形状是否可以处理，避免lxml错误
        
        Args:
            shape: 要检查的形状对象
            verbose: 是否输出详细信息
            
        Returns:
            bool: True表示形状安全可处理，False表示应该跳过
        """
        try:
            # 检查是否为组合形状
            if hasattr(shape, 'shape_type'):
                if shape.shape_type == 6:  # MSO_SHAPE_TYPE.GROUP
                    if verbose:
                        print(f"  跳过组合形状 (shape_type=6)")
                    return False
                # 可以在这里添加其他需要跳过的形状类型
                # 比如某些特殊的形状类型也可能导致问题
        except Exception as e:
            if verbose:
                print(f"  检查形状类型时出错，跳过该形状: {e}")
            return False
        
        # 尝试访问基本属性来验证形状的有效性
        try:
            # 尝试访问一些基本属性
            _ = hasattr(shape, 'has_text_frame')
            _ = hasattr(shape, 'has_table')
            return True
        except Exception as e:
            if verbose:
                print(f"  形状属性访问异常，跳过该形状: {e}")
            return False
    
    def _find_best_layout(self, presentation):
        """
        查找最佳的幻灯片布局
        """
        # 首先尝试找到名称包含"blank"的布局
        for layout in presentation.slide_layouts:
            if 'blank' in layout.name.lower():
                return layout
        
        # 如果没找到，尝试找到占位符数量最少的布局
        min_placeholders = float('inf')
        best_layout = None
        for layout in presentation.slide_layouts:
            placeholder_count = len(layout.placeholders)
            if placeholder_count < min_placeholders:
                min_placeholders = placeholder_count
                best_layout = layout
        
        # 如果仍然没有，使用第一个布局
        return best_layout or presentation.slide_layouts[0]
    
    def _copy_slide_content(self, source_slide, target_slide):
        """
        完整复制幻灯片内容，包括图片
        """
        try:
            # 清空目标幻灯片的所有内容，包括占位符
            for shape in list(target_slide.shapes):
                sp = shape.element
                sp.getparent().remove(sp)
            
            # 复制源幻灯片的所有形状（包括占位符内容）
            for shape in source_slide.shapes:
                # 安全检查形状
                if not self._is_safe_shape(shape, verbose=False):
                    continue
                    
                try:
                    # 复制形状的XML元素
                    from copy import deepcopy
                    shape_element = deepcopy(shape.element)
                    
                    # 添加到目标幻灯片
                    target_slide.shapes._spTree.append(shape_element)
                    
                    # 如果是图片，需要特殊处理
                    if hasattr(shape, 'shape_type'):
                        from pptx.enum.shapes import MSO_SHAPE_TYPE
                        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                            # 复制图片资源
                            self._copy_picture(shape, target_slide, shape_element)
                            
                except Exception as shape_error:
                    print(f"⚠️ 形状复制失败: {str(shape_error)}")
                    continue
                    
        except Exception as e:
            print(f"⚠️ 幻灯片内容复制失败: {str(e)}")
            raise
    
    def _copy_picture(self, source_shape, target_slide, target_shape_element):
        """
        复制图片资源
        """
        try:
            from pptx.enum.shapes import MSO_SHAPE_TYPE
            
            # 如果源形状是图片
            if hasattr(source_shape, 'shape_type') and source_shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                # 尝试重新添加图片
                try:
                    # 获取图片数据
                    image_part = source_shape.image.blob
                    
                    # 获取图片位置和大小
                    left = source_shape.left
                    top = source_shape.top  
                    width = source_shape.width
                    height = source_shape.height
                    
                    # 从目标幻灯片中移除复制的空形状
                    target_shape_element.getparent().remove(target_shape_element)
                    
                    # 重新添加图片到正确位置
                    from io import BytesIO
                    image_stream = BytesIO(image_part)
                    
                    picture = target_slide.shapes.add_picture(image_stream, left, top, width, height)
                    
                    print(f"🖼️ 图片重新添加成功")
                    
                except Exception as img_error:
                    print(f"⚠️ 图片重新添加失败: {str(img_error)}")
                    # 如果失败，保留原来复制的形状XML
                    pass
                        
        except Exception as e:
            print(f"⚠️ 图片资源复制失败: {str(e)}")
            pass
    
    def _copy_slide(self, source_slide, target_presentation):
        """
        完整的幻灯片复制方法
        
        Args:
            source_slide: 源幻灯片
            target_presentation: 目标演示文稿
            
        Returns:
            新创建的幻灯片对象
        """
        # 查找最佳的空白布局
        blank_layout = self._find_best_layout(target_presentation)
        new_slide = target_presentation.slides.add_slide(blank_layout)
        
        # 完整复制幻灯片内容
        self._copy_slide_content(source_slide, new_slide)
        
        return new_slide
    
    def _copy_presentation_settings(self, source_prs, target_prs, verbose: bool = True):
        """
        复制演示文稿的页面设置，包括幻灯片尺寸、比例等
        """
        try:
            # 获取源演示文稿的页面设置
            source_slide_width = source_prs.slide_width
            source_slide_height = source_prs.slide_height
            
            # 设置目标演示文稿的页面尺寸
            target_prs.slide_width = source_slide_width
            target_prs.slide_height = source_slide_height
            
            if verbose:
                print(f"📐 复制页面设置: {source_slide_width} x {source_slide_height}")
                
        except Exception as e:
            # 如果复制设置失败，使用默认设置
            if verbose:
                print(f"⚠️ 页面设置复制失败，使用默认设置: {str(e)}")
            pass
    
    def _render_slides_from_instructions(self, template_prs, render_instructions, data, verbose: bool = True):
        """
        根据渲染指令创建演示文稿
        
        Args:
            template_prs: 原始模板演示文稿
            render_instructions: 渲染指令列表，格式为 [(template_index, data_path), ...]
            data: 数据字典
            verbose: 是否显示详细信息
            
        Returns:
            新创建的演示文稿
        """
        if not render_instructions:
            return template_prs
        
        total_template_slides = len(template_prs.slides)
        
        if verbose:
            print(f"🎯 开始按渲染指令创建演示文稿")
            print(f"📊 模板幻灯片数量: {total_template_slides}")
            print(f"🔢 渲染指令数量: {len(render_instructions)}")
        
        # 验证渲染指令中的模板索引
        for i, (template_idx, data_path) in enumerate(render_instructions):
            if template_idx < 0 or template_idx >= total_template_slides:
                if verbose:
                    print(f"❌ 无效的模板索引: {template_idx} (模板只有 {total_template_slides} 张幻灯片)")
                return template_prs
        
        # 创建新的演示文稿
        result_prs = Presentation()
        self._copy_presentation_settings(template_prs, result_prs, verbose)
        
        # 删除默认幻灯片
        if len(result_prs.slides) > 0:
            slide_id = result_prs.slides[0].slide_id
            result_prs.part.drop_rel(result_prs.slides._sldIdLst[0].rId)
            del result_prs.slides._sldIdLst[0]
        
        # 按渲染指令复制和填充幻灯片
        for new_pos, (template_idx, data_path) in enumerate(render_instructions):
            # 复制模板幻灯片
            template_slide = template_prs.slides[template_idx]
            new_slide = self._copy_slide(template_slide, result_prs)
            
            # 获取该数据路径对应的数据
            slide_data = get_value_by_key(data, data_path)
            
            if verbose:
                print(f"📋 位置 {new_pos+1}: 使用模板 {template_idx+1}, 数据路径: {data_path}")
                if slide_data is not None:
                    print(f"   数据获取成功: {type(slide_data).__name__}")
                else:
                    print(f"   ⚠️ 数据路径无效或数据为空")
            
            # 填充当前幻灯片的数据
            self._fill_single_slide(new_slide, slide_data or {}, verbose)
        
        if verbose:
            print(f"✅ 渲染完成! 输出 {len(render_instructions)} 张幻灯片")
        
        return result_prs
    
    def _fill_single_slide(self, slide, slide_data, verbose: bool = True):
        """
        填充单个幻灯片的数据
        
        Args:
            slide: 幻灯片对象
            slide_data: 该幻灯片的数据
            verbose: 是否显示详细信息
        """
        # 如果slide_data不是字典，将其转换为字典格式
        if not isinstance(slide_data, dict):
            # 如果是简单值，创建一个通用的数据结构
            processed_slide_data = {
                "content": slide_data,
                "title": str(slide_data) if slide_data is not None else "",
                "description": str(slide_data) if slide_data is not None else ""
            }
        else:
            processed_slide_data = slide_data
        
        # 处理远程图片下载（复用原有逻辑）
        image_extensions = ('.jpg', '.jpeg', '.png', '.gif', '.bmp', '.webp', '.tiff', '.svg')
        
        def process_value(value):
            if isinstance(value, str):
                if value.endswith('.csv') and os.path.exists(value):
                    return convert_csv_to_json_list(value)
                elif value.startswith(('http://', 'https://')):
                    url_path = value.split('?')[0].lower()
                    if url_path.endswith(image_extensions):
                        local_image_path = download_image(value)
                        if local_image_path:
                            return local_image_path
                return value
            elif isinstance(value, list):
                return [process_value(item) for item in value]
            elif isinstance(value, dict):
                return {k: process_value(v) for k, v in value.items()}
            else:
                return value
        
        # 处理数据中的远程资源
        for key, value in processed_slide_data.items():
            processed_value = process_value(value)
            if processed_value is not None:
                processed_slide_data[key] = processed_value
        
        # 表格填充
        slide_table_requests = []
        slide_tables = []
        
        for shape in slide.shapes:
            # 安全检查形状
            if not self._is_safe_shape(shape, verbose):
                continue
            
            try:
                if shape.has_text_frame:
                    text = shape.text.strip()
                    if text.startswith("{{#") and text.endswith("}}"):
                        key = text[3:-2].strip()
                        table_data = get_value_by_key(processed_slide_data, key)
                        
                        if isinstance(table_data, str) and table_data.endswith('.csv') and os.path.exists(table_data):
                            table_data = convert_csv_to_json_list(table_data)
                        
                        if table_data is not None and isinstance(table_data, list):
                            slide_table_requests.append((shape, key, table_data))
                
                if hasattr(shape, 'has_table') and shape.has_table:
                    slide_tables.append(shape)
            except Exception as e:
                if verbose:
                    print(f"  处理形状时出错，跳过该形状: {e}")
                continue
        
        # 处理表格填充
        self._fill_slide_tables(slide, slide_table_requests, slide_tables)
        
        # 文本和图片填充
        for shape in list(slide.shapes):
            # 安全检查形状
            if not self._is_safe_shape(shape, verbose):
                continue
            
            try:
                if not shape.has_text_frame:
                    continue
            except Exception as e:
                if verbose:
                    print(f"  访问text_frame属性时出错，跳过该形状: {e}")
                continue
            
            try:
                text = shape.text.strip()
            except Exception as e:
                if verbose:
                    print(f"  访问text属性时出错，跳过该形状: {e}")
                continue
            
            if "{{" in text and "}}" in text:
                pure_placeholder = is_pure_placeholder(text)
                
                if pure_placeholder:
                    key = pure_placeholder
                    content_type = "text"

                    if key.startswith("@"):
                        key = key[1:]
                        content_type = "image"
                    elif key.startswith("#"):
                        continue

                    value = get_value_by_key(processed_slide_data, key)
                    if value is None:
                        continue

                    if content_type == "text":
                        if isinstance(value, str) and any(marker in value for marker in ['*', '#', '`', '\n']):
                            parse_markdown_text_preserve_format(shape.text_frame, value)
                        elif isinstance(value, list):
                            process_list_preserve_format(shape.text_frame, value)
                        else:
                            replace_text_preserve_format(shape.text_frame, str(value))

                    elif content_type == "image":
                        left, top, width, height = shape.left, shape.top, shape.width, shape.height
                        slide.shapes._spTree.remove(shape._element)
                            
                        if os.path.exists(value):
                            slide.shapes.add_picture(value, left, top, width=width, height=height)
                            if verbose:
                                print(f"✅ 成功替换图片: {key}")
                        elif value.startswith(('http://', 'https://')):
                            if verbose:
                                print(f"❌ 图片处理失败：远程图片下载失败 - {key}")
                            text_box = slide.shapes.add_textbox(left, top, width, height)
                            text_frame = text_box.text_frame
                            text_frame.text = f"图片加载失败: {key}"
                        else:
                            if verbose:
                                print(f"⚠️ 警告: 图片文件不存在: {value}")
                            text_box = slide.shapes.add_textbox(left, top, width, height)
                            text_frame = text_box.text_frame
                            text_frame.text = f"图片不存在: {key}"
                
                else:
                    # 混合文本模式
                    replaced_text = replace_mixed_placeholders(text, processed_slide_data)
                    replace_text_preserve_format(shape.text_frame, replaced_text)
    


    def _fill_slide_tables(self, slide, table_requests, slide_tables):
        """
        处理单个页面的表格填充，确保不会跨页面匹配
        
        Args:
            slide: 当前页面对象
            table_requests: 当前页面的表格占位符请求列表 [(shape, key, data)]
            slide_tables: 当前页面的表格列表
        """
        if not table_requests:
            return
        
        shapes_to_remove = []
        processed_tables = set()
        
        for placeholder_shape, key, table_data in table_requests:
            # 只在当前页面的表格中查找
            available_tables = [t for t in slide_tables if id(t) not in processed_tables]
            if not available_tables:
                # 如果当前页面的表格都被处理过，则允许重复使用
                available_tables = slide_tables
            
            if available_tables:
                nearest_table_shape = find_nearest_table(placeholder_shape, available_tables)
                if nearest_table_shape:
                    print(f"占位符 '{{#{key}}}' 匹配到当前页面的表格")
                    fill_existing_table(nearest_table_shape.table, table_data)
                    processed_tables.add(id(nearest_table_shape))
                else:
                    print(f"警告: 占位符 '{{#{key}}}' 在当前页面未找到合适的表格")
            else:
                print(f"警告: 当前页面没有可用的表格来填充占位符 '{{#{key}}}'")
            
            shapes_to_remove.append(placeholder_shape)
        
        # 删除当前页面的表格占位符文本框
        for shape in shapes_to_remove:
            try:
                shape._element.getparent().remove(shape._element)
            except Exception as e:
                print(f"删除占位符时出错: {e}")

    def fill(self, 
             data: dict, 
             template_file_path: str, 
             output_file_path: Optional[str] = None,
             output_format: str = "local",
             personal_auth_key: Optional[str] = None,
             personal_auth_secret: Optional[str] = None,
             base_url: str = "https://uat.agentspro.cn",
             order_info: Optional[List[tuple]] = None,
             verbose: bool = True) -> Union[str, Dict]:
        """
        使用嵌套JSON数据结构填充PowerPoint模板，支持动态重排序
        
        Args:
            data: 要填充的数据字典，支持嵌套结构和点号路径访问
                 - 文本占位符: {{key}} 或 {{nested.key}} 对应嵌套数据
                 - 图片占位符: {{@key}} 或 {{@nested.key}} (图片路径或URL)
                 - 表格占位符: {{#key}} 或 {{#nested.key}} (列表数据)
            template_file_path: 模板文件路径（支持本地路径和URL）
            output_file_path: 输出文件路径（当output_format为"local"时必需）
            output_format: 输出格式，支持 "local"、"base64"、"url"
            personal_auth_key: 个人认证密钥（当output_format为"url"时需要）
            personal_auth_secret: 个人认证密钥（当output_format为"url"时需要）
            base_url: 上传服务的基础URL
            order_info: 渲染指令列表，格式为 [(template_index, data_path), ...]
                每个元组包含：(模板索引, 数据路径)
                例如：[(0, "report_cover"), (1, "achievements[0]"), (1, "achievements[1]")]
            verbose: 是否显示详细输出信息
            
        Returns:
            str: 当output_format为"local"时返回文件路径，为"base64"时返回base64字符串
            Dict: 当output_format为"url"时返回上传结果字典
            
        Example:
            data = {
                "report_cover": {"title": "年度总结报告", "author": "战略发展部"},
                "achievements": [
                    {"id": "A1", "desc": "完成了A轮融资，金额超预期。"},
                    {"id": "A2", "desc": "核心产品用户数突破100万大关。"}
                ],
                "challenges": [
                    {"id": "C1", "desc": "市场竞争加剧，需要寻找新的突破口。"}
                ],
                "plan_next_year": {"focus": "重点投入AI驱动的新产品线研发。"}
            }
            
            # 渲染指令
            render_instructions = [
                (0, "report_cover"),      # 模板0, 封面
                (1, "achievements[0]"),   # 模板1, 成就页
                (1, "achievements[1]"),   # 模板1, 成就页  
                (2, "challenges[0]"),     # 模板2, 挑战页
                (3, "plan_next_year"),    # 模板3, 计划页
            ]
            
            模板中的占位符：
            - {{title}} -> "我的演示文稿"
            - {{user.nickname}} -> "frank"
            - {{user.dad.nickname}} -> "frank-dad"
            - {{user.hobbies}} -> ["读书", "运动", "音乐"] (作为列表显示)
            - {{user.hobbies[0]}} -> "读书" (数组第一项)
            - {{company.products[1].name}} -> "产品B" (嵌套数组索引)
            - {{@company.logo}} -> 替换为图片
            - {{#company.products}} -> 填充到最近的表格
            
            调用示例：
            agent = PPTX2PPTXAgent()
            
            result = agent.fill(
                data=data,
                template_file_path="template.pptx",
                output_file_path="output.pptx",
                order_info=render_instructions
            )
        """
        
        # 参数验证
        if output_format not in ["local", "base64", "url"]:
            raise ValueError(f"不支持的输出格式: {output_format}，支持的格式: local, base64, url")
        
        if output_format == "local" and not output_file_path:
            raise ValueError("当output_format为'local'时，必须提供output_file_path参数")
            
        if output_format == "url" and not personal_auth_key and not personal_auth_secret:
            raise ValueError("当output_format为'url'时，必须提供jwt_token参数")
        
        # 🔄 标准化数据格式，自动处理外层包装
        if verbose:
            print(f"📊 原始数据格式检查...")
        data = normalize_data_format(data)
        if verbose:
            print(f"✅ 数据格式标准化完成")
        
        # 用于存储需要清理的临时文件
        temp_files = []
        
        # 检查模板路径是否为URL，如果是则下载到临时文件
        actual_template_path = template_file_path
        is_template_from_url = False
        
        if template_file_path.startswith(('http://', 'https://')):
            if verbose:
                print(f"检测到URL模板: {template_file_path}")
            downloaded_template = download_template(template_file_path)
            if downloaded_template:
                actual_template_path = downloaded_template
                temp_files.append(downloaded_template)
                is_template_from_url = True
                if verbose:
                    print(f"模板下载成功: {downloaded_template}")
            else:
                raise ValueError(f"无法下载模板文件: {template_file_path}")
        
        # 加载 PPTX 模板
        prs = Presentation(actual_template_path)

        # 处理远程图片下载
        processed_data = {}
        # 支持的图片文件后缀
        image_extensions = ('.jpg', '.jpeg', '.png', '.gif', '.bmp', '.webp', '.tiff', '.svg')
        
        def process_value(value):
            """递归处理数据值，支持CSV文件读取"""
            if isinstance(value, str):
                # 检查是否是CSV文件路径
                if value.endswith('.csv') and os.path.exists(value):
                    print(f"检测到CSV文件: {value}")
                    return convert_csv_to_json_list(value)
                # 检查是否是远程图片URL
                elif value.startswith(('http://', 'https://')):
                    # 先检查URL是否以图片文件后缀结尾（忽略查询参数）
                    url_path = value.split('?')[0].lower()  # 去掉查询参数并转为小写
                    is_image_by_extension = url_path.endswith(image_extensions)
                    
                    # 如果URL不以图片后缀结尾，也尝试下载，通过Content-Type判断
                    if is_image_by_extension:
                        print(f"检测到图片URL（基于扩展名）: {value}")
                    else:
                        print(f"检测到可能的图片URL（需验证Content-Type）: {value}")
                    
                    # 尝试下载远程图片
                    local_image_path = download_image(value)
                    if local_image_path:
                        temp_files.append(local_image_path)
                        print(f"成功下载图片: {value} -> {local_image_path}")
                        return local_image_path
                    else:
                        print(f"⚠️ 图片下载失败，保留原始URL: {value}")
                        # 保留原始URL而不是返回None，这样可以在后续处理中看到问题
                        return value
                else:
                    return value
            elif isinstance(value, list):
                # 递归处理列表中的每个元素
                return [process_value(item) for item in value]
            elif isinstance(value, dict):
                # 递归处理字典中的每个值
                return {k: process_value(v) for k, v in value.items()}
            else:
                return value

        for key, value in data.items():
            processed_value = process_value(value)
            if processed_value is not None:
                processed_data[key] = processed_value

        # 如果没有渲染指令，使用传统的全量填充模式
        if not order_info:
            # 1. 表格填充 - 按页面分组处理，避免跨页面匹配
            if verbose:
                print(f"开始扫描PPT模板中的占位符...")
            
            for slide_idx, slide in enumerate(prs.slides):
                if verbose:
                    print(f"扫描第 {slide_idx + 1} 页...")
                
                # 收集当前页面的表格占位符和表格
                slide_table_requests = []  # 当前页面的表格占位符
                slide_tables = []  # 当前页面的表格
                
                for shape in slide.shapes:
                    # 安全检查形状
                    if not self._is_safe_shape(shape, verbose):
                        continue
                    
                    try:
                        # 收集表格占位符
                        if shape.has_text_frame:
                            text = shape.text.strip()
                            if text.startswith("{{") and text.endswith("}}"):
                                if verbose:
                                    print(f"  找到占位符: {text}")
                            if text.startswith("{{#") and text.endswith("}}"):
                                key = text[3:-2].strip()  # 去掉 {{# 和 }}
                                if verbose:
                                    print(f"找到表格占位符: {{#{key}}}")
                                table_data = get_value_by_key(processed_data, key)
                                
                                # 如果表格数据是字符串，可能是CSV文件路径，需要处理
                                if isinstance(table_data, str):
                                    if table_data.endswith('.csv') and os.path.exists(table_data):
                                        print(f"检测到CSV文件: {table_data}")
                                        table_data = convert_csv_to_json_list(table_data)
                                    else:
                                        print(f"表格数据为字符串但不是有效的CSV文件: {table_data}")
                                        table_data = None
                                
                                if table_data is not None and isinstance(table_data, list):
                                    print(f"表格占位符 {{#{key}}} 数据解析成功，{len(table_data)} 条记录")
                                    slide_table_requests.append((shape, key, table_data))
                                else:
                                    print(f"表格占位符 {{#{key}}} 数据解析失败或格式不正确")
                        
                        # 收集当前页面的表格
                        if hasattr(shape, 'has_table') and shape.has_table:
                            slide_tables.append(shape)
                    except Exception as e:
                        if verbose:
                            print(f"  处理形状时出错，跳过该形状: {e}")
                        continue
                
                # 处理当前页面的表格填充
                self._fill_slide_tables(slide, slide_table_requests, slide_tables)
            
            # 2. 文本、图片填充
            for slide in prs.slides:
                for shape in list(slide.shapes):  # list() to allow removal
                    # 安全检查形状
                    if not self._is_safe_shape(shape, verbose):
                        continue
                    
                    try:
                        if not shape.has_text_frame:
                            continue
                    except Exception as e:
                        if verbose:
                            print(f"  访问text_frame属性时出错，跳过该形状: {e}")
                        continue
                    
                    try:
                        text = shape.text.strip()
                    except Exception as e:
                        if verbose:
                            print(f"  访问text属性时出错，跳过该形状: {e}")
                        continue
                    
                    # 检查是否包含占位符
                    if "{{" in text and "}}" in text:
                        # 检查是否为纯占位符
                        pure_placeholder = is_pure_placeholder(text)
                        
                        if pure_placeholder:
                            # 纯占位符模式（原有逻辑）
                            key = pure_placeholder
                            content_type = "text"

                            # 判断类型前缀
                            if key.startswith("@"):
                                key = key[1:]
                                content_type = "image"
                            elif key.startswith("#"):
                                # 表格已经在上面处理过了，跳过
                                continue

                            value = get_value_by_key(processed_data, key)
                            if value is None:
                                continue

                            if content_type == "text":
                                # 检查是否包含Markdown格式
                                if isinstance(value, str) and any(marker in value for marker in ['*', '#', '`', '\n']):
                                    # 使用Markdown解析，现在会保留格式
                                    parse_markdown_text_preserve_format(shape.text_frame, value)
                                elif isinstance(value, list):
                                    # 处理列表数据，每项作为bullet point，完全保留格式
                                    process_list_preserve_format(shape.text_frame, value)
                                else:
                                    # 普通文本 - 使用新的格式保留函数
                                    replace_text_preserve_format(shape.text_frame, str(value))

                            elif content_type == "image":
                                # 获取位置并删除原文本框
                                left, top, width, height = shape.left, shape.top, shape.width, shape.height
                                slide.shapes._spTree.remove(shape._element)
                                    
                                # 检查是否为本地文件路径或URL
                                if os.path.exists(value):
                                    # 本地文件路径
                                    slide.shapes.add_picture(value, left, top, width=width, height=height)
                                    print(f"✅ 成功替换图片 (本地文件): {key}")
                                elif value.startswith(('http://', 'https://')):
                                    # URL路径（下载失败的情况）
                                    print(f"❌ 图片占位符 {{@{key}}} 处理失败：远程图片下载失败")
                                    print(f"   原始URL: {value}")
                                    # 可以选择添加一个错误提示文本框
                                    text_box = slide.shapes.add_textbox(left, top, width, height)
                                    text_frame = text_box.text_frame
                                    text_frame.text = f"图片加载失败: {key}"
                                else:
                                    # 其他情况
                                    print(f"⚠️ 警告: 图片文件不存在: {value}")
                                    # 添加错误提示文本框
                                    text_box = slide.shapes.add_textbox(left, top, width, height)
                                    text_frame = text_box.text_frame
                                    text_frame.text = f"图片不存在: {key}"
                        
                        else:
                            # 混合文本模式（新功能）- 完全保留格式
                            replaced_text = replace_mixed_placeholders(text, processed_data)
                            replace_text_preserve_format(shape.text_frame, replaced_text)
                            if verbose:
                                print(f"混合文本替换: '{text}' -> '{replaced_text}'")

        # 3. 渲染指令处理（如果提供了order_info）
        else:
            if verbose:
                print(f"\n🔄 检测到渲染指令，开始处理...")
            
            # 验证渲染指令格式
            if not isinstance(order_info, list):
                raise ValueError(f"渲染指令必须是列表格式，当前类型: {type(order_info)}")
            
            for i, item in enumerate(order_info):
                if not isinstance(item, tuple) or len(item) != 2:
                    raise ValueError(f"渲染指令第{i}项格式错误，必须是(模板索引, 数据路径)的元组格式")
            
            if verbose:
                print(f"✅ 渲染指令格式验证通过")
            
            # 使用渲染指令从模板生成演示文稿
            prs = self._render_slides_from_instructions(prs, order_info, processed_data, verbose)

        # 根据输出格式处理结果
        result = None
        temp_output_path = None
        
        try:
            if output_format == "local":
                # 直接保存到指定路径
                prs.save(output_file_path)
                if verbose:
                    print(f"✅ PPT已保存到: {output_file_path}")
                result = output_file_path
                
            elif output_format == "base64":
                # 保存到临时文件，然后转换为base64
                temp_fd, temp_output_path = tempfile.mkstemp(suffix='.pptx')
                os.close(temp_fd)
                temp_files.append(temp_output_path)
                
                prs.save(temp_output_path)
                
                # 读取文件并转换为base64
                with open(temp_output_path, 'rb') as f:
                    file_bytes = f.read()
                    base64_str = base64.b64encode(file_bytes).decode('utf-8')
                
                if verbose:
                    print(f"✅ PPT已转换为base64格式 (大小: {len(base64_str)} 字符)")
                result = base64_str
                
            elif output_format == "url":
                # 保存到临时文件，然后上传
                temp_fd, temp_output_path = tempfile.mkstemp(suffix='.pptx')
                os.close(temp_fd)
                temp_files.append(temp_output_path)
                
                prs.save(temp_output_path)
                
                # 创建上传器并上传文件
                uploader = SimpleFileUploader(personal_auth_key, personal_auth_secret, base_url)
                
                # 生成文件名
                filename = f"filled_presentation_{os.path.basename(temp_output_path)}"
                
                with open(temp_output_path, 'rb') as f:
                    upload_result = uploader.upload(f, filename)
                
                if upload_result.get("success"):
                    if verbose:
                        print(f"✅ PPT已上传成功，文件ID: {base_url}/api/fs/{upload_result['fileId']}")
                    result = {
                        "fileId": upload_result['fileId'],
                        "fileUrl": f"{base_url}/api/fs/{upload_result['fileId']}"
                    }
                else:
                    raise Exception(f"文件上传失败: {upload_result.get('error', '未知错误')}")
            
        finally:
            # 清理临时文件（包括下载的图片、模板文件和输出临时文件）
            for temp_file in temp_files:
                try:
                    cleanup_temp_file(temp_file)
                    if verbose:
                        print(f"清理临时文件: {temp_file}")
                except Exception as e:
                    if verbose:
                        print(f"清理临时文件失败: {temp_file}, 错误: {e}")
        
        return result
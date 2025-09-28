from ..client import ChatClient

class HTML2PPTXAgent:
    def __init__(self):
        # 用于存储从封面页学习到的风格
        self.learned_style = None
        
        # 添加调试模式
        self.debug = True
    
    def _extract_style_from_html(self, html_content: str) -> dict:
        """从生成的HTML中提取样式信息"""
        style_info = {
            "colors": [],
            "fonts": [],
            "background": "",
            "layout_type": "",
            "design_elements": []
        }
        
        try:
            import re
            
            # 提取颜色信息
            color_patterns = [
                r'color:\s*([^;]+)',
                r'background:\s*([^;]+)',
                r'background-color:\s*([^;]+)',
                r'#[0-9a-fA-F]{3,6}'
            ]
            
            for pattern in color_patterns:
                matches = re.findall(pattern, html_content)
                style_info["colors"].extend(matches)
            
            # 提取字体信息
            font_matches = re.findall(r'font-family:\s*([^;]+)', html_content)
            style_info["fonts"].extend(font_matches)
            
            # 提取背景信息
            bg_matches = re.findall(r'background(?:-image)?:\s*([^;]+)', html_content)
            if bg_matches:
                style_info["background"] = bg_matches[0]
            
            # 判断布局类型
            if 'gradient' in html_content.lower():
                style_info["layout_type"] = "gradient_background"
            if 'card' in html_content.lower() or 'container' in html_content.lower():
                style_info["layout_type"] += "_card_layout"
            
            # 提取设计元素
            if 'border-radius' in html_content:
                style_info["design_elements"].append("rounded_corners")
            if 'box-shadow' in html_content:
                style_info["design_elements"].append("shadows")
            if 'transform' in html_content:
                style_info["design_elements"].append("transforms")
                
        except Exception as e:
            print(f"样式提取出错: {e}")
        
        return style_info
    
    def _create_style_continuation_prompt(self, content_type: str, content: str) -> str:
        """创建延续封面风格的提示词（已注释，让agent自主识别布局）"""
        # 注释掉风格延续机制，让agent为不同页面类型选择最适合的布局
        # if self.learned_style:
        #     style_desc = f"""
        # 请延续封面页的设计风格生成{content_type}：
        # 
        # 参考风格信息：
        # - 主要颜色：{', '.join(self.learned_style.get('colors', [])[:5])}
        # - 字体样式：{', '.join(self.learned_style.get('fonts', [])[:2])}
        # - 背景风格：{self.learned_style.get('background', '渐变背景')}
        # - 布局类型：{self.learned_style.get('layout_type', '卡片式布局')}
        # - 设计元素：{', '.join(self.learned_style.get('design_elements', []))}
        # 
        # 内容：{content}
        # 
        # 请保持与封面页一致的视觉风格，生成{content_type}的HTML代码。"""
        # else:
        #     # 如果还没有学习到风格，让agent自主选择
        #     style_desc = f"""
        # 内容：{content}
        # 
        # 请根据内容自动识别合适的设计风格，生成{content_type}的HTML代码。"""
        # 
        # return style_desc
        
        # 直接返回内容，让agent自主识别页面类型和布局
        return content

    def _create_style_prompt(self, content_type: str, content: str) -> str:
        """创建包含样式指导的提示词（已注释，让agent自主识别）"""
        # 注释掉预设风格，让agent自主识别内容类型并选择合适的布局
        # style_guide = f"""
        # 设计要求：
        # - 主题：沈从文《边城》文学分析
        # - 色调：{self.design_theme['color_scheme']}
        # - 风格：{self.design_theme['style']}
        # - 布局：{self.design_theme['layout']}
        # - 字体：{self.design_theme['font']}
        # - 背景：{self.design_theme['background']}
        # 
        # 页面类型：{content_type}
        # 内容：{content}
        # 
        # 请生成一致风格的HTML代码。"""
        # return style_guide
        
        # 直接返回内容，让agent自主识别和处理
        return content
    
    def _debug_print(self, title: str, content: str):
        """调试输出函数"""
        if self.debug:
            print(f"\n{'='*20} {title} {'='*20}")
            print(content)
            print(f"{'='*50}")

    def _clean_html_content(self, html_content: str) -> str:
        """清理HTML内容，去除markdown代码块标记和重复内容，只保留HTML文档部分"""
        cleaned_html = html_content.strip()
        
        # 去除markdown代码块标记
        if cleaned_html.startswith('```html'):
            cleaned_html = cleaned_html[7:]  # 去除开头的```html
        if cleaned_html.endswith('```'):
            cleaned_html = cleaned_html[:-3]  # 去除结尾的```
        cleaned_html = cleaned_html.strip()
        
        # 查找HTML文档的开始和结束位置
        doctype_start = cleaned_html.find('<!DOCTYPE html>')
        html_end = cleaned_html.rfind('</html>')
        
        # 如果找到完整的HTML文档，只提取这部分
        if doctype_start != -1 and html_end != -1:
            cleaned_html = cleaned_html[doctype_start:html_end + 7]  # +7是为了包含</html>
        
        # 处理重复的HTML内容，只保留第一个完整的HTML文档
        if cleaned_html.count('<!DOCTYPE html>') > 1:
            # 找到第一个</html>的位置
            first_html_end = cleaned_html.find('</html>')
            if first_html_end != -1:
                cleaned_html = cleaned_html[:first_html_end + 7]  # +7是为了包含</html>
        
        return cleaned_html

    def cover(self, outline_content: str):
        # 尝试不同的分割方式来提取封面内容
        if "### 封面页" in outline_content:
            cover_content = outline_content.split("### 封面页")[1].split("###")[0].strip()
        elif "**幻灯片1：封面页**" in outline_content:
            cover_content = outline_content.split("**幻灯片1：封面页**")[1].split("**幻灯片2")[0].strip()
        else:
            # 如果找不到特定格式，就使用前几行作为封面内容
            lines = outline_content.split('\n')[:10]
            cover_content = '\n'.join(lines)
        
        print("=== COVER CONTENT ===")
        print(cover_content)
        print("=== END COVER CONTENT ===")
        
        # 调试输出
        self._debug_print("封面页提取内容", cover_content)
        
        chat_client = ChatClient(
            agent_id="34e8c6cfe6ee448ea5191561a19d822e",
            personal_auth_key="48cf18e0e0ca4b51bbf8fa60193ffb5c",
            personal_auth_secret="HWlQXZ5vxgrXDGEtTGGdsTFhJfr9rCmD",
            base_url="https://uat.agentspro.cn"
        )

        # 使用简化的提示词，让agent自主识别
        prompt = self._create_style_prompt("封面页", cover_content)
        
        html_content = ""
        for event in chat_client.invoke(prompt):
            if event['type'] == 'start_bubble':
                print(f"\n{'=' * 20} 封面生成开始 {'=' * 20}")
            elif event['type'] == 'token':
                print(event['content'], end='', flush=True)
                html_content += event['content']
            elif event['type'] == 'end_bubble':
                print(f"\n{'=' * 20} 封面生成结束 {'=' * 20}")
            elif event['type'] == 'finish':
                print(f"\n{'=' * 20} 封面对话完成 {'=' * 20}")
                break
        
        # 预处理HTML内容，去除markdown代码块标记和重复内容
        cleaned_html = self._clean_html_content(html_content)
        
        # 暂时注释掉风格学习机制，观察agent自主识别效果
        # self.learned_style = self._extract_style_from_html(cleaned_html)
        # self._debug_print("学习到的风格信息", str(self.learned_style))
        
        # 保存为HTML文件
        with open('cover.html', 'w', encoding='utf-8') as f:
            f.write(cleaned_html)
        print(f"\n封面页已保存为 cover.html")
        print(f"Agent自主识别封面页类型和风格")

    def catalog(self, outline_content: str):
        # 尝试不同的分割方式来提取目录内容
        if "### 目录页" in outline_content:
            catalog_content = outline_content.split("### 目录页")[1].split("###")[0].strip()
        elif "**幻灯片2：目录页**" in outline_content:
            catalog_content = outline_content.split("**幻灯片2：目录页**")[1].split("**幻灯片3")[0].strip()
        else:
            # 如果找不到特定格式，尝试提取目录相关内容
            lines = outline_content.split('\n')
            catalog_lines = []
            for line in lines:
                if any(keyword in line for keyword in ["引言", "一、", "二、", "三、", "结语", "总结"]):
                    catalog_lines.append(line.strip())
            catalog_content = '\n'.join(catalog_lines[:10])  # 取前10个目录项
        
        print("=== CATALOG CONTENT ===")
        print(catalog_content)
        print("=== END CATALOG CONTENT ===")
        
        # 调试输出
        self._debug_print("目录页提取内容", catalog_content)
        
        chat_client = ChatClient(
            agent_id="34e8c6cfe6ee448ea5191561a19d822e",
            personal_auth_key="48cf18e0e0ca4b51bbf8fa60193ffb5c",
            personal_auth_secret="HWlQXZ5vxgrXDGEtTGGdsTFhJfr9rCmD",
            base_url="https://uat.agentspro.cn"
        )
        
        # 使用简化的提示词，让agent自主识别
        prompt = self._create_style_prompt("目录页", catalog_content)
        
        html_content = ""
        for event in chat_client.invoke(prompt):
            if event['type'] == 'start_bubble':
                print(f"\n{'=' * 20} 目录生成开始 {'=' * 20}")
            elif event['type'] == 'token':
                print(event['content'], end='', flush=True)
                html_content += event['content']
            elif event['type'] == 'end_bubble':
                print(f"\n{'=' * 20} 目录生成结束 {'=' * 20}")
            elif event['type'] == 'finish':
                print(f"\n{'=' * 20} 目录对话完成 {'=' * 20}")
                break
        
        # 预处理HTML内容，去除markdown代码块标记和重复内容
        cleaned_html = self._clean_html_content(html_content)
        
        # 保存为HTML文件
        with open('catalog.html', 'w', encoding='utf-8') as f:
            f.write(cleaned_html)
        print(f"\n目录页已保存为 catalog.html")
        print("Agent自主识别目录页布局和风格")

    def conclusion(self, outline_content: str):
        # 尝试不同的分割方式来提取结论内容
        if "### 总结页" in outline_content:
            conclusion_content = outline_content.split("### 总结页")[1].strip()
        elif "**总结页**" in outline_content:
            conclusion_content = outline_content.split("**总结页**")[1].strip()
        elif "结语" in outline_content:
            # 提取结语部分
            lines = outline_content.split('\n')
            conclusion_lines = []
            start_collecting = False
            for line in lines:
                if "结语" in line or "总结" in line:
                    start_collecting = True
                if start_collecting:
                    conclusion_lines.append(line.strip())
                    if len(conclusion_lines) > 15:  # 限制行数
                        break
            conclusion_content = '\n'.join(conclusion_lines)
        else:
            # 如果找不到特定格式，使用最后几行作为结论内容
            lines = outline_content.split('\n')
            conclusion_content = '\n'.join(lines[-10:])
        
        print("=== CONCLUSION CONTENT ===")
        print(conclusion_content)
        print("=== END CONCLUSION CONTENT ===")
        
        # 调试输出
        self._debug_print("结论页提取内容", conclusion_content)
        
        chat_client = ChatClient(
            agent_id="34e8c6cfe6ee448ea5191561a19d822e",
            personal_auth_key="48cf18e0e0ca4b51bbf8fa60193ffb5c",
            personal_auth_secret="HWlQXZ5vxgrXDGEtTGGdsTFhJfr9rCmD",
            base_url="https://uat.agentspro.cn"
        )
        
        # 使用简化的提示词，让agent自主识别
        prompt = self._create_style_prompt("结论页", conclusion_content)
        
        html_content = ""
        for event in chat_client.invoke(prompt):
            if event['type'] == 'start_bubble':
                print(f"\n{'=' * 20} 结论生成开始 {'=' * 20}")
            elif event['type'] == 'token':
                print(event['content'], end='', flush=True)
                html_content += event['content']
            elif event['type'] == 'end_bubble':
                print(f"\n{'=' * 20} 结论生成结束 {'=' * 20}")
            elif event['type'] == 'finish':
                print(f"\n{'=' * 20} 结论对话完成 {'=' * 20}")
                break
        
        # 预处理HTML内容，去除markdown代码块标记和重复内容
        cleaned_html = self._clean_html_content(html_content)
        
        # 保存为HTML文件
        with open('conclusion.html', 'w', encoding='utf-8') as f:
            f.write(cleaned_html)
        print(f"\n结论页已保存为 conclusion.html")
        print("Agent自主识别结论页布局和风格")
    
    def content(self, outline_content: str):
        """生成内容页"""
        # 尝试不同的分割方式来提取内容页内容
        if "### 内容页" in outline_content:
            content_content = outline_content.split("### 内容页")[1].split("###")[0].strip()
        elif "**幻灯片3：内容页**" in outline_content:
            content_content = outline_content.split("**幻灯片3：内容页**")[1].split("**幻灯片4")[0].strip()
        elif "**幻灯片4：内容页**" in outline_content:
            content_content = outline_content.split("**幻灯片4：内容页**")[1].split("**幻灯片5")[0].strip()
        else:
            # 如果找不到特定格式，尝试提取主要内容部分
            lines = outline_content.split('\n')
            content_lines = []
            start_extracting = False
            for line in lines:
                # 跳过封面和目录相关内容，开始提取主体内容
                if any(keyword in line for keyword in ["内容", "正文", "主体", "核心", "技术", "方案", "产品"]):
                    start_extracting = True
                if start_extracting and line.strip():
                    content_lines.append(line.strip())
                    if len(content_lines) >= 15:  # 提取足够的内容
                        break
            content_content = '\n'.join(content_lines)
        
        print("=== CONTENT CONTENT ===")
        print(content_content)
        print("=== END CONTENT CONTENT ===")
        
        # 调试输出
        self._debug_print("内容页提取内容", content_content)
        
        chat_client = ChatClient(
            agent_id="34e8c6cfe6ee448ea5191561a19d822e",
            personal_auth_key="48cf18e0e0ca4b51bbf8fa60193ffb5c",
            personal_auth_secret="HWlQXZ5vxgrXDGEtTGGdsTFhJfr9rCmD",
            base_url="https://uat.agentspro.cn"
        )
        
        # 使用简化的提示词，让agent自主识别
        prompt = self._create_style_prompt("内容页", content_content)
        
        html_content = ""
        for event in chat_client.invoke(prompt):
            if event['type'] == 'start_bubble':
                print(f"\n{'=' * 20} 内容页生成开始 {'=' * 20}")
            elif event['type'] == 'token':
                print(event['content'], end='', flush=True)
                html_content += event['content']
            elif event['type'] == 'end_bubble':
                print(f"\n{'=' * 20} 内容页生成结束 {'=' * 20}")
            elif event['type'] == 'finish':
                print(f"\n{'=' * 20} 内容页对话完成 {'=' * 20}")
                break
        
        # 预处理HTML内容，去除markdown代码块标记和重复内容
        cleaned_html = self._clean_html_content(html_content)
        
        # 保存为HTML文件
        with open('content.html', 'w', encoding='utf-8') as f:
            f.write(cleaned_html)
        print(f"\n内容页已保存为 content.html")
        print("Agent自主识别内容页布局和风格")
    
    def debug_content_extraction(self, outline_content: str):
        """调试内容提取，帮助诊断问题"""
        print("\n" + "="*60)
        print("                  调试信息                  ")
        print("="*60)
        
        print(f"\n1. 原始outline_content长度: {len(outline_content)}")
        print(f"2. 原始outline_content前500字符:")
        print("-" * 40)
        print(outline_content[:500])
        print("-" * 40)
        
        # 测试封面内容提取
        if "### 封面页" in outline_content:
            cover_content = outline_content.split("### 封面页")[1].split("###")[0].strip()
            extraction_method = "### 封面页格式"
        elif "**幻灯片1：封面页**" in outline_content:
            cover_content = outline_content.split("**幻灯片1：封面页**")[1].split("**幻灯片2")[0].strip()
            extraction_method = "**幻灯片1：封面页**格式"
        else:
            lines = outline_content.split('\n')[:10]
            cover_content = '\n'.join(lines)
            extraction_method = "默认前10行"
        
        print(f"\n3. 封面内容提取方式: {extraction_method}")
        print(f"4. 提取的封面内容:")
        print("-" * 40)
        print(cover_content)
        print("-" * 40)
        
        # 测试目录内容提取
        lines = outline_content.split('\n')
        catalog_lines = []
        for line in lines:
            if any(keyword in line for keyword in ["引言", "一、", "二、", "三、", "结语", "总结"]):
                catalog_lines.append(line.strip())
        catalog_content = '\n'.join(catalog_lines[:10])
        
        print(f"\n5. 提取的目录内容:")
        print("-" * 40)
        print(catalog_content)
        print("-" * 40)
        
        # 生成的样式提示词示例
        style_prompt = self._create_style_prompt("封面页", cover_content)
        print(f"\n6. 发送给agent的封面页提示词:")
        print("-" * 40)
        print(style_prompt)
        print("-" * 40)
        
        print(f"\n7. Agent自主识别机制说明:")
        print("-" * 40)
        print("• 封面页：agent根据内容自动识别类型并选择合适风格")
        print("• 目录页：agent自主选择最适合目录展示的布局")
        print("• 结论页：agent自主选择最适合总结的布局风格")
        print("• 不预设任何风格限制，完全由agent智能判断")
        print("• 测试agent泛化能力和自适应性")
        print("-" * 40)
        
        print("\n" + "="*60)
        print("                调试信息结束                ")
        print("="*60)

    def save(self, file_path: str):
        """将生成的HTML文件整合成PPTX"""
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.enum.text import PP_ALIGN
            from pptx.dml.color import RGBColor
            import os
            import re
            from bs4 import BeautifulSoup
        except ImportError as e:
            print(f"缺少必要的库: {e}")
            print("请安装: pip install python-pptx beautifulsoup4")
            return
        
        # 创建新的演示文稿
        prs = Presentation()
        
        # 定义要处理的HTML文件列表（按顺序）
        html_files = [
            ('cover.html', '封面页'),
            ('catalog.html', '目录页'), 
            ('content.html', '内容页'),
            ('conclusion.html', '结论页')
        ]
        
        print(f"开始将HTML文件整合为PPTX: {file_path}")
        
        for html_file, page_type in html_files:
            if os.path.exists(html_file):
                print(f"正在处理 {html_file} ({page_type})...")
                
                # 读取HTML内容
                with open(html_file, 'r', encoding='utf-8') as f:
                    html_content = f.read()
                
                # 解析HTML内容
                soup = BeautifulSoup(html_content, 'html.parser')
                
                # 添加新幻灯片（使用空白布局）
                slide_layout = prs.slide_layouts[6]  # 空白布局
                slide = prs.slides.add_slide(slide_layout)
                
                # 提取标题
                title_text = self._extract_title_from_html(soup)
                if title_text:
                    # 添加标题文本框
                    title_box = slide.shapes.add_textbox(
                        Inches(0.5), Inches(0.5), Inches(9), Inches(1.5)
                    )
                    title_frame = title_box.text_frame
                    title_frame.text = title_text
                    title_paragraph = title_frame.paragraphs[0]
                    title_paragraph.alignment = PP_ALIGN.CENTER
                    title_font = title_paragraph.runs[0].font
                    title_font.size = Pt(36)
                    title_font.bold = True
                    title_font.color.rgb = RGBColor(0, 51, 102)  # 深蓝色
                
                # 提取主要内容
                content_text = self._extract_content_from_html(soup, page_type)
                if content_text:
                    # 添加内容文本框
                    content_box = slide.shapes.add_textbox(
                        Inches(0.5), Inches(2.5), Inches(9), Inches(5)
                    )
                    content_frame = content_box.text_frame
                    content_frame.text = content_text
                    content_frame.word_wrap = True
                    
                    # 设置内容字体
                    for paragraph in content_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.size = Pt(18)
                            run.font.color.rgb = RGBColor(51, 51, 51)
                
                print(f"  {page_type} 添加完成")
            else:
                print(f"警告: {html_file} 文件不存在，跳过...")
        
        # 保存PPTX文件
        prs.save(file_path)
        print(f"\nPPTX文件已保存为: {file_path}")
        print(f"总共生成了 {len(prs.slides)} 张幻灯片")
    
    def _extract_title_from_html(self, soup):
        """从HTML中提取标题"""
        # 尝试多种方式提取标题
        title_selectors = [
            'h1',
            '.main-title',
            '.header-title', 
            'title',
            '.title'
        ]
        
        for selector in title_selectors:
            title_elem = soup.select_one(selector)
            if title_elem:
                return title_elem.get_text().strip()
        
        # 如果找不到标题，返回None
        return None
    
    def _extract_content_from_html(self, soup, page_type):
        """从HTML中提取主要内容"""
        content_parts = []
        
        if page_type == '目录页':
            # 提取目录项
            list_items = soup.find_all(['li', 'div'])
            for item in list_items:
                text = item.get_text().strip()
                if text and len(text) > 3:  # 过滤掉太短的文本
                    content_parts.append(f"• {text}")
        
        elif page_type == '内容页':
            # 提取段落和列表
            content_elements = soup.find_all(['p', 'li', 'div'])
            for elem in content_elements:
                text = elem.get_text().strip()
                if text and len(text) > 10:  # 过滤掉太短的文本
                    content_parts.append(text)
        
        elif page_type == '结论页':
            # 提取结论要点
            content_elements = soup.find_all(['p', 'li', 'div', 'span'])
            for elem in content_elements:
                text = elem.get_text().strip()
                if text and len(text) > 5:
                    content_parts.append(text)
        
        else:  # 封面页
            # 提取副标题和描述
            subtitle_selectors = [
                '.subtitle',
                '.subtitle-chinese',
                '.subtitle-english',
                'p',
                '.description'
            ]
            for selector in subtitle_selectors:
                elems = soup.select(selector)
                for elem in elems:
                    text = elem.get_text().strip()
                    if text and len(text) > 5:
                        content_parts.append(text)
        
        # 限制内容长度，避免幻灯片过于拥挤
        if len(content_parts) > 8:
            content_parts = content_parts[:8]
            content_parts.append("...")
        
        return '\n\n'.join(content_parts)
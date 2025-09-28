"""
PowerPoint Presentation Parser Module

This module provides functionality for parsing Microsoft PowerPoint presentations
(.ppt, .pptx, .pptm) into structured representations. It can extract text content,
process embedded images, and organize the presentation by slides.
"""

import os
import subprocess
import tempfile
from collections.abc import MutableSequence
from pathlib import Path
from typing import Literal, cast

from rsb.models.field import Field

from agentle.generations.models.message_parts.file import FilePart
from agentle.generations.models.structured_outputs_store.visual_media_description import (
    VisualMediaDescription,
)
from agentle.generations.providers.base.generation_provider import GenerationProvider
from agentle.parsing.document_parser import DocumentParser
from agentle.parsing.image import Image
from agentle.parsing.parsed_file import ParsedFile
from agentle.parsing.section_content import SectionContent


class PptxFileParser(DocumentParser):
    """
    Parser for processing Microsoft PowerPoint presentations (.ppt, .pptx, .pptm).

    This parser extracts content from PowerPoint presentations, including text and embedded
    images. Each slide in the presentation is parsed as a separate section in the resulting
    ParsedFile. With the "high" strategy, embedded images are analyzed using a visual
    description agent to extract text via OCR and generate descriptions.

    The parser supports both modern PowerPoint formats (.pptx, .pptm) and the legacy
    format (.ppt), converting the latter to .pptx using LibreOffice when necessary.

    **Attributes:**

    *   `strategy` (Literal["high", "low"]):
        The parsing strategy to use. Defaults to "high".
        - "high": Performs thorough parsing including OCR and image analysis
        - "low": Performs basic text extraction without analyzing images

        **Example:**
        ```python
        parser = PptxFileParser(strategy="low")  # Use faster, less intensive parsing
        ```

    *   `visual_description_agent` (Agent[VisualMediaDescription]):
        The agent used to analyze and describe image content. If provided and
        strategy is "high", this agent will be used to analyze images embedded
        in the presentation.
        Defaults to the agent created by `visual_description_agent_default_factory()`.

        **Example:**
        ```python
        from agentle.agents.agent import Agent
        from agentle.generations.models.structured_outputs_store.visual_media_description import VisualMediaDescription

        custom_agent = Agent(
            model="gemini-2.0-pro-vision",
            instructions="Focus on diagram and chart analysis in presentations",
            response_schema=VisualMediaDescription
        )

        parser = PptxFileParser(visual_description_agent=custom_agent)
        ```

    *   `multi_modal_provider` (GenerationProvider):
        An alternative to using a visual_description_agent. This is a generation
        provider capable of handling multi-modal content (text and images).
        Defaults to GoogleGenerationProvider().

        Note: You cannot use both visual_description_agent and multi_modal_provider
        at the same time.

    **Usage Examples:**

    Basic parsing of a PowerPoint presentation:
    ```python
    from agentle.parsing.parsers.pptx import PptxFileParser

    # Create a parser with default settings
    parser = PptxFileParser()

    # Parse a PowerPoint file
    parsed_presentation = parser.parse("presentation.pptx")

    # Access the slides (as sections)
    for i, section in enumerate(parsed_presentation.sections):
        print(f"Slide {i+1} content:")
        print(section.text[:100] + "...")  # Print first 100 chars of each slide
    ```

    Processing images in a presentation:
    ```python
    from agentle.parsing.parsers.pptx import PptxFileParser

    # Create a parser with high-detail strategy
    parser = PptxFileParser(strategy="high")

    # Parse a presentation with images
    slide_deck = parser.parse("slide_deck.pptx")

    # Extract and process images
    for i, section in enumerate(slide_deck.sections):
        slide_num = i + 1
        print(f"Slide {slide_num} has {len(section.images)} images")

        for j, image in enumerate(section.images):
            print(f"  Image {j+1}:")
            if image.ocr_text:
                print(f"    OCR text: {image.ocr_text}")
    ```

    **Requirements:**

    For parsing .ppt (legacy format) files, LibreOffice must be installed on the system
    to perform the conversion to .pptx. If LibreOffice is not installed, a RuntimeError
    will be raised when attempting to parse .ppt files.
    """

    type: Literal["pptx"] = "pptx"

    strategy: Literal["high", "low"] = Field(default="high")

    visual_description_provider: GenerationProvider | None = Field(
        default=None,
    )
    """
    The agent to use for generating the visual description of the document.
    Useful when you want to customize the prompt for the visual description.
    """

    async def parse_async(
        self,
        document_path: str,
    ) -> ParsedFile:
        """
        Asynchronously parse a PowerPoint presentation and generate a structured representation.

        This method reads a PowerPoint file, extracts text and image content from each slide,
        and processes embedded images to extract text and generate descriptions when using
        the "high" strategy.

        For .ppt (legacy format) files, the method first converts them to .pptx format
        using LibreOffice before processing.

        Args:
            document_path (str): Path to the PowerPoint file to be parsed

        Returns:
            ParsedFile: A structured representation where:
                - Each slide is a separate section
                - Text content is extracted from each slide
                - Images are extracted and (optionally) analyzed

        Raises:
            RuntimeError: If converting a .ppt file fails (e.g., if LibreOffice is not installed)
                or if the conversion times out

        Example:
            ```python
            import asyncio
            from agentle.parsing.parsers.pptx import PptxFileParser

            async def process_presentation():
                parser = PptxFileParser(strategy="high")
                result = await parser.parse_async("slide_deck.pptx")

                # Print information about the slides
                print(f"Presentation has {len(result.sections)} slides")

                # Access slide content
                for i, section in enumerate(result.sections):
                    print(f"Slide {i+1} content:")
                    print(section.text[:100] + "...")

                    # Check for images
                    if section.images:
                        print(f"  Contains {len(section.images)} images")

            asyncio.run(process_presentation())
            ```

        Note:
            This method uses the python-pptx library to read PowerPoint files. For optimal
            results with legacy .ppt files, ensure LibreOffice is installed on the system.
        """
        import hashlib

        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        from pptx.presentation import Presentation as PptxPresentation
        from pptx.shapes.autoshape import Shape
        from pptx.shapes.picture import Picture

        path = Path(document_path)
        converted_pptx_file: str | None = None
        if path.suffix in {".ppt", ".pptm"}:
            converted_pptx_file = self._convert_to_pptx(document_path)

        prs: PptxPresentation = Presentation(converted_pptx_file or document_path)
        sections: MutableSequence[SectionContent] = []
        processed_images: dict[str, tuple[str, str]] = {}

        for slide_index, slide in enumerate(prs.slides, start=1):
            slide_texts: list[str] = []
            slide_images: list[tuple[str, bytes, str]] = []  # (name, data, hash)

            for shape in slide.shapes:
                if shape.has_text_frame:
                    shape_with_text = cast(Shape, shape)
                    text_str: str = shape_with_text.text
                    slide_texts.append(text_str)

                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    picture_shape = cast(Picture, shape)
                    image_blob: bytes = picture_shape.image.blob
                    image_hash = hashlib.sha256(image_blob).hexdigest()
                    image_name: str = (
                        shape.name or f"slide_{slide_index}_img_{image_hash[:8]}"
                    )
                    slide_images.append((image_name, image_blob, image_hash))

            combined_text: str = "\n".join(slide_texts)
            final_images: MutableSequence[Image] = []
            image_descriptions: MutableSequence[str] = []

            if self.visual_description_provider and self.strategy == "high":
                for img_idx, (image_name, image_blob, image_hash) in enumerate(
                    slide_images, start=1
                ):
                    is_cached = image_hash in processed_images
                    if is_cached:
                        cached_md, cached_ocr = processed_images[image_hash]
                        image_descriptions.append(
                            f"Slide {slide_index} - Image {img_idx}: {cached_md}"
                        )
                        final_images.append(
                            Image(
                                name=image_name,
                                contents=image_blob,
                                ocr_text=cached_ocr,
                            )
                        )
                        continue

                    agent_input = FilePart(
                        mime_type=Path(image_name).suffix,
                        data=image_blob,
                    )
                    agent_response = await self.visual_description_provider.generate_by_prompt_async(
                        agent_input,
                        developer_prompt="You are a helpful assistant that deeply understands visual media.",
                        response_schema=VisualMediaDescription,
                    )
                    image_md: str = agent_response.parsed.md
                    image_ocr = agent_response.parsed.ocr_text

                    processed_images[image_hash] = (image_md, image_ocr or "")
                    image_descriptions.append(
                        f"Slide {slide_index} - Image {img_idx}: {image_md}"
                    )
                    final_images.append(
                        Image(name=image_name, contents=image_blob, ocr_text=image_ocr)
                    )

                if image_descriptions:
                    combined_text += "\n\n" + "\n".join(image_descriptions)

            section_content = SectionContent(
                number=slide_index,
                text=combined_text,
                md=combined_text,
                images=final_images,
            )
            sections.append(section_content)

        return ParsedFile(
            name=path.name,
            sections=sections,
        )

    def _convert_to_pptx(self, document_path: str) -> str:
        """
        Convert a legacy PowerPoint file (.ppt or .pptm) to the modern .pptx format.

        This helper method uses LibreOffice to convert legacy PowerPoint files to
        the modern .pptx format, which can then be processed by the python-pptx library.

        Args:
            document_path (str): Path to the legacy PowerPoint file (.ppt or .pptm)

        Returns:
            str: Path to the converted .pptx file

        Raises:
            RuntimeError: If LibreOffice is not installed, the conversion fails,
                or the conversion times out

        Note:
            This method requires LibreOffice to be installed on the system and
            available in the PATH. The conversion is performed in a temporary
            directory and the converted file is returned as a temporary file
            that will be automatically cleaned up when no longer needed.
        """

        def _is_libreoffice_installed() -> bool:
            try:
                subprocess.run(
                    ["libreoffice", "--version"],
                    stdout=subprocess.DEVNULL,
                    stderr=subprocess.DEVNULL,
                    check=True,
                )
                return True
            except (subprocess.CalledProcessError, FileNotFoundError):
                return False

        if not _is_libreoffice_installed():
            raise RuntimeError("LibreOffice not found in system PATH")

        with tempfile.TemporaryDirectory() as temp_dir:
            # Write input file to temporary directory
            input_filename = Path(document_path).name
            input_path = os.path.join(temp_dir, input_filename)
            with (
                open(document_path, "rb") as src_file,
                open(input_path, "wb") as dst_file,
            ):
                dst_file.write(src_file.read())

            # Run LibreOffice conversion
            try:
                subprocess.run(
                    [
                        "libreoffice",
                        "--headless",
                        "--convert-to",
                        "pptx",
                        "--outdir",
                        temp_dir,
                        input_path,
                    ],
                    check=True,
                    stdout=subprocess.DEVNULL,
                    stderr=subprocess.PIPE,
                    timeout=60,
                )
            except subprocess.CalledProcessError as e:
                error_msg = e.stderr.decode().strip() if e.stderr else "Unknown error"
                raise RuntimeError(f"Conversion failed: {error_msg}") from e
            except subprocess.TimeoutExpired:
                raise RuntimeError("Conversion timed out after 60 seconds")

            # Determine output file path
            output_filename = Path(input_filename).stem + ".pptx"
            output_path = os.path.join(temp_dir, output_filename)

            if not os.path.exists(output_path):
                available_files = os.listdir(temp_dir)
                raise RuntimeError(
                    f"Converted file not found at {output_path}. Found files: {available_files}"
                )

            # Read the converted file and create a temporary file in the system temp directory
            with open(output_path, "rb") as f:
                content = f.read()

            temp_output = tempfile.NamedTemporaryFile(suffix=".pptx", delete=False)
            temp_output.write(content)
            temp_output.close()

            return temp_output.name
